# -*- coding: utf-8 -*-
import requests
import pdfplumber
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import os
import re
import google.generativeai as genai
from dotenv import load_dotenv
from transformers import pipeline
from moviepy.editor import VideoFileClip, ImageClip, concatenate_videoclips, CompositeVideoClip, TextClip, ColorClip, AudioFileClip
from moviepy.editor import *
from moviepy.video.fx.all import crop, loop
from moviepy.video.tools.drawing import circle
import json
import pysrt
import numpy as np
from gtts import gTTS
from googletrans import Translator
from pydub import AudioSegment
from moviepy.video.fx.all import colorx
from moviepy.video.tools.subtitles import SubtitlesClip
import base64
from docx import Document
from pptx import Presentation
import time
import logging
from requests.exceptions import RequestException

# Load environment variables
load_dotenv()

# Constants
pictures_folder = "pictures"
videos_folder = "videos"
output_video_path = "final_slideshow.mp4"
audio_output_path = "final_audio.mp3"
audio_output_speedup_path = "final_audio_speeded_up.mp3"
srt_file_path = "subtitles.srt"
output_folder = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\pdf_images"

# Create folders if they don't exist
os.makedirs(pictures_folder, exist_ok=True)
os.makedirs(videos_folder, exist_ok=True)
os.makedirs(output_folder, exist_ok=True)

#initialize sarvam api key
sarvam_api_key = os.getenv("SARVAM_API_KEY")

# Initialize Gemini API
api_key = os.getenv("API_KEY")
genai.configure(api_key=api_key)
model = genai.GenerativeModel('gemini-1.5-flash')

# Set path for Tesseract OCR
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Extract images from PDF
def extract_images_from_pdf(pdf_path, output_folder):
    doc = fitz.open(pdf_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_path = os.path.join(output_folder, f"page_{page_num}_img_{img_index}.png")
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
            yield image_path  # Yield each saved image path for OCR processing

# OCR on images in folder
def ocr_images_in_folder(folder_path):
    text = ""
    for filename in os.listdir(folder_path):
        if filename.endswith(".png"):
            image_path = os.path.join(folder_path, filename)
            try:
                image = Image.open(image_path)
                text += pytesseract.image_to_string(image) + "\n"
            except Exception as e:
                print(f"Error OCR processing {filename}: {e}")
    return text

# Extract text from images and save them
def extract_text_from_pdf_images(pdf_path, output_folder, output_path):
    print("extract_text_from_pdf_images TRIGGERED")
    text = ""
    for image_path in extract_images_from_pdf(pdf_path, output_folder):
        try:
            image = Image.open(image_path)
            text += pytesseract.image_to_string(image) + "\n"
        except Exception as e:
            print(f"Error OCR processing {image_path}: {e}")

    # Save OCR text from images to file
    with open(output_path, "a", encoding="utf-8") as file:  # Append to include both PDF text and image text
        file.write(text)
    print(f"OCR text from images extracted and saved to {output_path}")

# Extract text from txt file
def extract_text_from_txt(txt_path):
    with open(txt_path, 'r', encoding='utf-8') as file:
        return file.read()

# Extract text from docx file
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs])

# Extract text from pptx file
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Extract text from pdf file
def extract_text_from_pdf(pdf_path, output_path):
    print("extract_text_from_pdf TRIGGERED")
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    
    # Save extracted text to file
    with open(output_path, "w", encoding="utf-8") as file:
        file.write(text)
    return text

# Function to handle all file types
def extract_text_from_file(file_path, output_path):
    if file_path.endswith('.pdf'):
        # Extract text from both PDF content and images
        text = extract_text_from_pdf(file_path, output_path)
        extract_text_from_pdf_images(file_path, output_folder, output_path)  # Append OCR text from images
    elif file_path.endswith('.txt'):
        text = extract_text_from_txt(file_path)
    elif file_path.endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file format")
    
    # Write text to output file if not already saved within specific functions
    if text:
        with open(output_path, "w", encoding="utf-8") as file:
            file.write(text)
        print(f"Text extracted and saved to {output_path}")

# paths to the pdf file and the output text file
pdf_path = r"C:\Users\Happy yadav\Desktop\aura.ai\test\doc\budget_speech.pdf"
output_text_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\extracted.txt"
cleaned_text_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\cleaned.txt"

extract_text_from_file(pdf_path, output_text_path)

# Clean the extracted text
def clean_text(text):
    """Cleans the given text by removing extra spaces and special characters."""
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with a single space
    text = re.sub(r'[^a-zA-Z0-9\s.,]', '', text)  # Remove unwanted characters
    return text.strip()  # Remove leading and trailing spaces

# Remove UIN and CIN numbers from the text
def remove_uin_cin(text):
    # Pattern for UIN codes, assuming they start with 'UIN' followed by an alphanumeric string
    text = re.sub(r'\bUIN\s+[A-Z0-9]+\b', '', text)

    # Pattern for CIN codes, assuming they start with 'CIN' followed by an alphanumeric string
    text = re.sub(r'\bCIN\s+[A-Z0-9]+\b', '', text)

    # Remove extra spaces left after deletion
    text = re.sub(r'\s+', ' ', text).strip()
    return text

# SingleFunction calling the above functions
def clean_extracted_text(input_file, output_file):
    """Cleans the text from the input file and saves it to the output file."""
    with open(input_file, 'r', encoding='utf-8') as file:
        text = file.read()

    # Clean the text
    cleaned_text = clean_text(text)
    cleaned_text = remove_uin_cin(cleaned_text)

    # Save the cleaned text to a new file
    with open(output_file, 'w', encoding='utf-8') as file:
        file.write(cleaned_text)
    
    print(f"Cleaned text saved to '{output_file}'")

# save the cleaned text to a new file
clean_extracted_text(output_text_path, cleaned_text_path)

#if numerical data is present in the text, we can remove it using the below function

#insert the function if required

# calling the gemini api ffor script calling
def generate_script(cleaned_txt_path, video_length, language, llm_prompt_path):

    # Read cleaned text from the specified path
    with open(cleaned_txt_path, "r", encoding="utf-8") as file:
        cleaned_text = file.read()

    # Prepare the prompt
    llm_prompt = f"""
        <prompt>
            <step1>Extract the key information from the document and identify the main points that need to be discussed.</step1>
            <step2>Break down these main points into smaller subtopics that can be explained sequentially.</step2>
            <step3>For each subtopic, construct a simple and engaging explanation that fits naturally into a continuous narrative.</step3>
            <step4>Ensure that each subtopic flows logically into the next, maintaining coherence and a clear structure.</step4>
            <step5>Use simple and concise language, keeping the audience's attention with relatable examples or anecdotes where necessary.</step5>
            <step6>Review the entire monologue to make sure it includes all the important information from the original document.</step6>
            <step7>Don't include anything related to QR code or any link.</step7>
            <step8>Return the script only. The response should start with the script and end with it.</step8>
            <step9>Response should not contain anything like Script:</step9>
        </prompt>

        <!-- Text from cleaned.txt for reference -->
        {cleaned_text}
        
        <!-- Please make the script approximately {video_length} seconds long. -->
    """
    
    # Add language to the prompt if it's not English
    if language != "English":
        llm_prompt += f"\n\n<!-- Please make the script in {language}. -->"

    # Generate the content with Gemini API
    response = model.generate_content([llm_prompt])

    # Save the response to 'llm_prompt.txt'
    with open(llm_prompt_path, "w", encoding="utf-8") as output_file:
        output_file.write(response.text)
    print("Script generated successfully and saved to '{llm_prompt_path}'")
    
    return response.text

video_length = 40  # Example: 120 seconds for a 2-minute video
language = "English"  # Specify the language, e.g., "English", "Spanish", etc.
llm_prompt_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\llm_prompt.txt"

# Call the function and get the generated script
script_text = generate_script(cleaned_text_path, video_length, language, llm_prompt_path)

# Questions generation from the text
def generate_quiz_from_text(cleaned_text):
    quiz_string = ""
    try:
        # Define the prompt to generate quiz questions with different types, difficulties, and explanations
        llm_prompt = f"""
        You are provided with the following text: 

        {cleaned_text}

        Based on this text, generate a pool of 10 questions that include a variety of question types:
        - Multiple-choice (MCQ)
        - True/false
        - Fill in the blanks (with options)

        For each question, output the response in the following JSON format:

        [
            {{
            "question": "Question 1 text here?",
            "options": ["Option 1", "Option 2", "Option 3", "Option 4"],
            "correctAnswer": "Correct option here",
            "explanation": "Explanation for the answer here",
            "type": "mcq" or "true-false" or "fill-in-the-blank",
            "difficulty": "Easy" or "Medium" or "Hard",
            "use": "normal"
            }},
            {{
            "question": "Question 2 text here?",
            "options": ["Option 1", "Option 2", "Option 3", "Option 4"],
            "correctAnswer": "Correct option here",
            "explanation": "Explanation for the answer here",
            "type": "mcq" or "true-false" or "fill-in-the-blank",
            "difficulty": "Easy" or "Medium" or "Hard",
            "use": "substitute"
            }},
            ...
        ]

        Ensure:
        - The correct answers are accurate and based on the provided text.
        - The questions are a mix of different types (multiple-choice, true/false, fill-in-the-blanks).
        - Each question is labeled with the appropriate difficulty level: "Easy," "Medium," or "Hard". 
        - In total, there must be 20 questions. All types: true-false, MCQs, and fill-in-the-blanks (with 3 options) must be present.
        - I need 5 questions for the quiz and 5 substitutes of only easy and medium type. 
        - Provide detailed explanations for each correct answer.

        Reply with just the JSON response and nothing else.
        """

        # Generate content using the model
        response = model.generate_content([llm_prompt])
        quiz_string = response.text
    except Exception as e:
        print(f"Error generating quiz: {e}")
    
    return quiz_string

def save_quiz_to_json(quiz_string, output_file):
    # Clean up any code block formatting from the response
    cleaned_string = quiz_string.replace('```json', '').replace('```', '').strip()

    try:
        # Parse the cleaned string as JSON
        quiz_json = json.loads(cleaned_string)
        
        # Save to specified output file
        with open(output_file, 'w', encoding='utf-8') as json_file:
            json.dump(quiz_json, json_file, indent=4)
        
        print(f"Quiz data saved successfully to '{questions_file_path}'")
    except json.JSONDecodeError as e:
        print(f"Failed to decode JSON: {e}")

questions_file_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\questions.json"


# Load cleaned text
with open(llm_prompt_path, "r", encoding="utf-8") as file:
    cleaned_text = file.read()

# Generate quiz questions and save to JSON
quiz_string = generate_quiz_from_text(cleaned_text)
save_quiz_to_json(quiz_string, questions_file_path)

def chunk_text(text, max_length=500):
    """Splits text into chunks of a maximum length."""
    words = text.split()
    chunks = []
    current_chunk = ""
    
    for word in words:
        if len(current_chunk) + len(word) + 1 <= max_length:
            current_chunk += (word + " ")
        else:
            chunks.append(current_chunk.strip())
            current_chunk = word + " "
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks

def generate_audio_with_background(llm_prompt_path, output_audio_path, background_music_path, final_output_path, language_code="en-IN", speaker="meera", pitch=0.1, pace=1, loudness=1, music_volume=-15):
    try:
        # Read text from llm_prompt_path
        with open(llm_prompt_path, "r", encoding="utf-8") as file:
            text_content = file.read()
        
        # Split text into chunks within the 500 character limit
        text_chunks = chunk_text(text_content)
        
        # List to hold audio segments
        audio_segments = []

        # Process each chunk with the Sarvam API
        url = "https://api.sarvam.ai/text-to-speech"
        headers = {
            "api-subscription-key": sarvam_api_key,
            "Content-Type": "application/json"
        }
        
        for chunk in text_chunks:
            payload = {
                "target_language_code": language_code,
                "speaker": speaker,
                "pitch": pitch,
                "pace": pace,
                "loudness": loudness,
                "enable_preprocessing": True,
                "model": "bulbul:v1",
                "inputs": [chunk]
            }
            
            response = requests.post(url, json=payload, headers=headers)
            
            # Check if the response is successful
            if response.status_code == 200:
                audio_data = response.json().get('audios')[0]
                audio_bytes = base64.b64decode(audio_data)
                
                # Load audio segment and append to list
                with open("temp_chunk.wav", "wb") as audio_file:
                    audio_file.write(audio_bytes)
                
                audio_segment = AudioSegment.from_file("temp_chunk.wav")
                audio_segments.append(audio_segment)
                os.remove("temp_chunk.wav")  # Clean up temp file
            else:
                print(f"Error processing chunk: {response.status_code} - {response.text}")
                return
        
        # Concatenate all audio segments
        main_audio = sum(audio_segments)
        main_audio.export(output_audio_path, format="wav")
        print(f"Concatenated audio saved to {output_audio_path}")

        # Process and mix audio with background music
        # Load the main audio
        main_audio = AudioSegment.from_file(output_audio_path)
        
        # Load the background music and adjust volume
        background_music = AudioSegment.from_file(background_music_path)
        background_music = background_music - abs(music_volume)  # Lower the volume

        # Loop background music to match the length of the main audio
        if len(background_music) < len(main_audio):
            loop_count = len(main_audio) // len(background_music) + 1
            background_music = background_music * loop_count
        background_music = background_music[:len(main_audio)]
        
        # Overlay main audio on background music
        final_audio = main_audio.overlay(background_music)

        # Export the final mixed audio
        final_audio.export(final_output_path, format="mp3")
        print(f"Final audio with background music saved to {final_output_path}")

        # Delete the intermediate audio file
        os.remove(output_audio_path)
        print(f"Intermediate file '{output_audio_path}' deleted successfully.")

    except Exception as e:
        print(f"Error generating audio with background music: {e}")

# Static Paths for the function
output_audio_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\sarvam_generated_audio.wav"
background_music_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\background.mp3"
final_output_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\sarvam_output_audio_music.mp3"

# Call the function to generate audio with background music using sarvam api
generate_audio_with_background(
    llm_prompt_path=llm_prompt_path,
    output_audio_path=output_audio_path,
    background_music_path=background_music_path,
    final_output_path=final_output_path,
    language_code="en-IN",
    speaker="meera",
    pitch=0.1,
    pace=1,
    loudness=1,
    music_volume=-15  # Adjust background music volume
)

# Funtion for Aura bot to answer the questions
def get_text_from_txt(file_path):
    """Reads text from a given file path."""
    try:
        with open(file_path, "r", encoding="utf-8") as text_file:
            text = text_file.read()
        return text
    except FileNotFoundError:
        print(f"The file {file_path} does not exist.")
        return None
    except Exception as e:
        print(f"An error occurred while reading the file: {e}")
        return None

# Function to ask a question to the Gemini AI model based on the llm_prompt file as knowledge base
def ask_aura_question(user_question, txt_file_path, model):
    """Answers the user's question using text from cleaned.txt or provides a general response."""
    answer = ""
    try:
        # Retrieve the extracted text from the txt file
        extracted_text = get_text_from_txt(txt_file_path)
        
        if not extracted_text:
            return "Error: Could not retrieve text from file."
        
        # Create a prompt asking Gemini to answer based on `cleaned.txt`
        prompt = f"Based on the following text, try to answer the question:\n\n'{extracted_text}'\n\nQuestion: {user_question}"
        response = model.generate_content(prompt)
        print("Initial answer response: ", response.text)  # Debug print for checking response
        answer = response.text

        # Check if the answer is relevant to the knowledge base
        if "I don't know" in answer or "I'm not sure" in answer or len(answer.strip()) < 5:
            # Provide a fallback response from Gemini's general knowledge if the response is inadequate
            fallback_prompt = f"I don't have specific knowledge about this part of your query based on the provided text, but I can answer from general knowledge.\n\nQuestion: {user_question}"
            fallback_response = model.generate_content(fallback_prompt)
            print("Fallback answer response: ", fallback_response.text)  # Debug print for checking fallback response
            answer = "I don’t have detailed information from the provided knowledge base, but here’s what I can tell you: " + fallback_response.text
    except Exception as e:
        print(f"Error while asking Gemini AI: {e}")
        return f"Error: {e}"
    
    return answer

response = ask_aura_question("What is the name of the policy ?", cleaned_text_path, model)
print("Answer:", response)

# Function to generate a script JSON file with timed subtitles
def generate_script_json(audio_path, llm_prompt_path, output_script_path, max_words_per_line=8):
    # Load the audio file to get the duration
    audio = AudioFileClip(audio_path)
    audio_duration = audio.duration  # in seconds

    # Read subtitle text from llm_prompt_path
    try:
        with open(llm_prompt_path, "r", encoding="utf-8") as file:
            subtitle_text = file.read()
    except FileNotFoundError:
        print(f"Error: The file {llm_prompt_path} does not exist.")
        return
    except Exception as e:
        print(f"An error occurred while reading the file: {e}")
        return

    # Split subtitle text into chunks
    words = subtitle_text.split()
    subtitle_chunks = []
    current_chunk = ""
    
    for word in words:
        # Account for max words per line in each chunk
        if len(current_chunk.split()) < max_words_per_line:
            current_chunk += f"{word} "
        else:
            subtitle_chunks.append(current_chunk.strip())
            current_chunk = f"{word} "
    subtitle_chunks.append(current_chunk.strip())  # Add the last chunk

    # Calculate duration for each chunk, taking small pauses for commas
    chunk_durations = []
    remaining_duration = audio_duration
    for chunk in subtitle_chunks:
        # Give a small additional pause if commas are present
        comma_pause = chunk.count(",") * 0.03  # 0.05s pause for each comma
        duration = (audio_duration / len(subtitle_chunks)) + comma_pause
        chunk_durations.append(min(duration, remaining_duration))
        remaining_duration -= duration
    
    # Calculate start and end times for each chunk
    start_time = 0.0
    script_data = []

    for i, (chunk, duration) in enumerate(zip(subtitle_chunks, chunk_durations)):
        end_time = start_time + duration
        script_data.append({
            "id": i + 1,
            "start_time": round(start_time, 2),
            "end_time": round(end_time, 2),
            "chunk": chunk
        })
        start_time = end_time  # Set start time for next chunk
    
    # Save the script data to a JSON file
    with open(output_script_path, "w", encoding="utf-8") as script_file:
        json.dump(script_data, script_file, indent=4, ensure_ascii=False)
    print(f"Script JSON file saved as {output_script_path}")

# Usage example
output_script_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\script.json"

generate_script_json(final_output_path, llm_prompt_path, output_script_path)

# Assume this function uses the Gemini API to generate content from a provided prompt
def generate_keywords_batch(text_chunks: list) -> list:
    # Define the batch prompt
    llm_prompt = """
    <prompt>
        For each chunk provided, generate a single visual metaphor keyword phrase of 4-5 words maximum that captures the main theme.The keyword should avoid ambiguous or complex terms and focus on general, easily interpretable metaphor visuals for example-(man in hospital, car accident). Only return the keyword as a single phrase without any bullet points, explanations, or formatting. The output should simply be one keyword per line.
    </prompt>
    """
    
    # Combine all chunks for batch processing
    batch_prompt = llm_prompt + "\n".join([f"Chunk {i+1}: {chunk}" for i, chunk in enumerate(text_chunks)])
    response = model.generate_content(batch_prompt)
    
    # Parse response to extract keywords
    keywords = [line.strip() for line in response.text.splitlines() if line.strip()]
    
    return keywords

# Secondary function to generate keywords based on a single prompt from llm_prompt_path content
def generate_missing_keywords(prompt_path: str, num_keywords: int) -> list:
    # Read the content of the prompt file
    with open(prompt_path, 'r') as f:
        prompt_content = f.read()
    
    # Create a prompt for missing keywords
    missing_prompt = f"""
    <prompt>
        Generate {num_keywords} unique visual metaphor keyword phrases, like (man in hospital, girl reading book, a car accident etc) each 4-5 words, capturing themes within this content. The keyword should avoid ambiguous terms and focus on easily interpretable, general visuals related to the content in that content. Only return the keyword as a single phrase without any bullet points, explanations, or formatting. The output should simply be one keyword per line. Also note that the keyword should be unique and not repeated, because these keywords will be used to generate the videos from pixaby api.
    </prompt>
    """
    response = model.generate_content(missing_prompt + prompt_content)
    
    # Parse the generated keywords
    missing_keywords = [line.strip() for line in response.text.splitlines() if line.strip()]
    
    return missing_keywords

# Main function to update script.json with keywords, ensuring no keywords are missing
def update_script_json_batch(script_file_path, llm_prompt_path):
    """
    Updates script.json with keywords using batch processing, checking for any missing keywords,
    and filling missing positions using a single Gemini API request based on llm_prompt_path content.
    """
    # Load the existing script.json
    with open(script_file_path, 'r') as f:
        script_data = json.load(f)
    
    # Process chunks in batches for efficiency
    batch_size = 10
    for i in range(0, len(script_data), batch_size):
        batch = script_data[i:i+batch_size]
        batch_texts = [chunk['chunk'] for chunk in batch]

        # Generate initial keywords for the batch
        keywords = generate_keywords_batch(batch_texts)
        
        # Assign keywords to each chunk in the batch
        for chunk, keyword in zip(batch, keywords):
            chunk['keyword'] = keyword

        # Sleep briefly to respect rate limits, if needed
        time.sleep(1)

    # Check for missing keywords in the updated script_data
    missing_indices = [i for i, chunk in enumerate(script_data) if not chunk.get('keyword')]
    num_missing_keywords = len(missing_indices)
    print(f"Number of missing keywords Found: {num_missing_keywords}")
    # If there are missing keywords, use llm_prompt_path content to generate them
    if num_missing_keywords > 0:
        # Generate missing keywords using Gemini API and llm_prompt_path content
        missing_keywords = generate_missing_keywords(llm_prompt_path, num_missing_keywords)
        
        # Place generated keywords in missing positions
        for i, missing_index in enumerate(missing_indices):
            script_data[missing_index]['keyword'] = missing_keywords[i]

    # Write the updated data back to script.json
    with open(script_file_path, 'w') as f:
        json.dump(script_data, f, indent=4)
    print("Script.json Ready!!")

# Example call to update_script_json_batch
update_script_json_batch(output_script_path, llm_prompt_path)

# Function to calculate video durations based on start and end times in script.json
def calculate_video_durations(script_file_path):
    """
    Adds video_duration to each chunk in script.json based on start_time and end_time.
    """
    # Load the script JSON data
    with open(script_file_path, 'r') as f:
        script_data = json.load(f)

    # Calculate video duration for each chunk and update the data
    for chunk in script_data:
        start_time = chunk["start_time"]
        end_time = chunk["end_time"]
        chunk["video_duration"] = round(end_time - start_time, 2)  # Calculate and round to 2 decimal places

    # Save the updated script.json with video durations
    with open(script_file_path, 'w') as f:
        json.dump(script_data, f, indent=4)
    print("Video durations added to script.json")

# Call this function before fetching or trimming videos
calculate_video_durations(output_script_path)

# logging :

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('image_generation.log'),
        logging.StreamHandler()
    ]
)

def download_image(prompt, image_path, width=1280, height=720, seed=468605, model='flux', max_retries=3, delay=5):
    """
    Downloads an image from Pollinations.ai with retry logic and error handling.
    """
    image_url = f"https://pollinations.ai/p/{prompt}?width={width}&height={height}&seed={seed}&model={model}"
    
    for attempt in range(max_retries):
        try:
            logging.info(f"Attempting to download image for prompt: {prompt} (Attempt {attempt + 1}/{max_retries})")
            response = requests.get(image_url, timeout=30)
            response.raise_for_status()
            
            with open(image_path, 'wb') as file:
                file.write(response.content)
            
            logging.info(f"Successfully saved image to {image_path}")
            return True
            
        except RequestException as e:
            logging.error(f"Attempt {attempt + 1} failed: {str(e)}")
            if attempt < max_retries - 1:
                logging.info(f"Retrying in {delay} seconds...")
                time.sleep(delay)
            continue
    
    logging.error(f"Failed to download image after {max_retries} attempts")
    return False

def generate_images_for_chunks(script_file_path):
    """
    Generates images for script chunks with error handling and logging.
    """
    try:
        with open(script_file_path, 'r') as f:
            script_data = json.load(f)
    except (json.JSONDecodeError, FileNotFoundError) as e:
        logging.error(f"Failed to load script.json: {str(e)}")
        return False

    images_dir = os.path.join(os.path.dirname(script_file_path), 'images')
    try:
        os.makedirs(images_dir, exist_ok=True)
    except OSError as e:
        logging.error(f"Failed to create images directory: {str(e)}")
        return False

    successful_downloads = 0
    failed_downloads = 0

    for chunk in script_data:
        prompt = chunk['keyword']
        image_filename = f"img_{chunk['id']}.jpg"
        image_path = os.path.join(images_dir, image_filename)
        
        if download_image(prompt, image_path):
            chunk['img_path'] = image_path
            successful_downloads += 1
        else:
            failed_downloads += 1
            chunk['img_path'] = None
        
        time.sleep(1)

    try:
        with open(script_file_path, 'w') as f:
            json.dump(script_data, f, indent=4)
        logging.info(f"Script.json updated. Successful downloads: {successful_downloads}, Failed: {failed_downloads}")
        return True
    except IOError as e:
        logging.error(f"Failed to update script.json: {str(e)}")
        return False
# ----------------------------------------------------------- working code for generating images using pollinations.ai
# def download_image(prompt, image_path, width=1280, height=720, seed=468605, model='flux'):
#     """
#     Downloads an image generated based on the given prompt from Pollinations.ai and saves it to the specified path.
#     """
#     # Format the image URL with the prompt and parameters
#     image_url = f"https://pollinations.ai/p/{prompt}?width={width}&height={height}&seed={seed}&model={model}"
    
#     # Request the image from Pollinations.ai
#     response = requests.get(image_url)
#     response.raise_for_status()  # Raise an error if the request was unsuccessful
    
#     # Save the image to the specified path
#     with open(image_path, 'wb') as file:
#         file.write(response.content)
    
#     print(f"Image downloaded and saved as {image_path}")

# def generate_images_for_chunks(script_file_path):
#     """
#     Iterates over chunks in script.json, generates an image for each chunk using Pollinations.ai,
#     saves each image in an 'images' directory, and updates script.json with the image file path.
#     """
#     # Load the script data from script.json
#     with open(script_file_path, 'r') as f:
#         script_data = json.load(f)
    
#     # Ensure 'images' directory exists
#     images_dir = os.path.join(os.path.dirname(script_file_path), 'images')
#     os.makedirs(images_dir, exist_ok=True)

#     # Process each chunk in the script.json
#     for chunk in script_data:
#         prompt = chunk['keyword']  # Use the chunk text as the prompt
#         image_filename = f"img_{chunk['id']}.jpg"
#         image_path = os.path.join(images_dir, image_filename)
        
#         # Generate and download the image
#         download_image(prompt, image_path)
        
#         # Update the chunk with the image path
#         chunk['img_path'] = image_path
        
#         # Sleep briefly to avoid overwhelming the API
#         time.sleep(1)
    
#     # Save the updated script.json with image paths
#     with open(script_file_path, 'w') as f:
#         json.dump(script_data, f, indent=4)
    
#     print("Script.json updated with image paths")

# ----------------------working CODE FOR GENERATING THE IMAGE USING STABILITY AI 
# def download_image(prompt, image_path, api_key, output_format="webp", aspect_ratio="16:9"):
#     """
#     Downloads an image generated based on the given prompt from Stability AI's API and saves it to the specified path.
#     """
#     url = "https://api.stability.ai/v2beta/stable-image/generate/core"
#     headers = {
#         "authorization": f"Bearer {api_key}",
#         "accept": "image/*"
#     }
#     data = {
#         "prompt": prompt,
#         "output_format": output_format,
#         "aspect_ratio": aspect_ratio
#     }
    
#     response = requests.post(url, headers=headers, data=data, files={"none": ""})
    
#     if response.status_code == 200:
#         with open(image_path, "wb") as file:
#             file.write(response.content)
#         print(f"Image downloaded and saved as {image_path}")
#     else:
#         raise Exception(f"Error: {response.status_code} - {response.text}")

# def generate_images_for_chunks(script_file_path, api_key, aspect_ratio="16:9"):
#     """
#     Iterates over chunks in script.json, generates an image for each chunk using Stability AI,
#     saves each image in an 'images' directory, and updates script.json with the image file path.
#     """
#     with open(script_file_path, "r") as f:
#         script_data = json.load(f)
    
#     images_dir = os.path.join(os.path.dirname(script_file_path), "images")
#     os.makedirs(images_dir, exist_ok=True)
    
#     for chunk in script_data:
#         prompt = chunk.get("keyword", "A generic AI-generated image")  # Default prompt if missing
#         image_filename = f"img_{chunk['id']}.webp"
#         image_path = os.path.join(images_dir, image_filename)
        
#         try:
#             download_image(prompt, image_path, api_key, aspect_ratio=aspect_ratio)
#             chunk["img_path"] = image_path
#         except Exception as e:
#             print(f"Failed to generate image for chunk {chunk['id']}: {e}")
        
#         time.sleep(1)  # Avoid hitting rate limits
    
#     with open(script_file_path, "w") as f:
#         json.dump(script_data, f, indent=4)
    
#     print("Script.json updated with image paths")

# Example usage (replace with your actual API key)
# api_key = "sk-Xips2HbQGXa2GTbRE3Y08s933bqp7eNhSUGw4eyFW1Vu4Cep"
# generate_images_for_chunks(output_script_path, api_key, aspect_ratio="16:9")


generate_images_for_chunks(output_script_path)

def generate_video_with_images_and_subtitles(audio_path, script_json_path, output_video_path, font_size=20):
    # Load audio file
    audio = AudioFileClip(audio_path)
    audio_duration = audio.duration

    # Load script.json
    try:
        with open(script_json_path, "r", encoding="utf-8") as json_file:
            script_data = json.load(json_file)
    except FileNotFoundError:
        print(f"Error: The file {script_json_path} does not exist.")
        return
    except Exception as e:
        print(f"An error occurred while reading the JSON file: {e}")
        return

    # Create a list to hold image clips with duration
    image_clips = []
    for entry in script_data:
        img_path = entry.get("img_path")
        duration = entry.get("video_duration")
        
        if not img_path or not os.path.exists(img_path):
            print(f"Image file not found for chunk: {entry['chunk']}")
            continue
        
        # Load the image and set its duration
        image_clip = ImageClip(img_path).set_duration(duration)
        image_clips.append(image_clip)
    
    # Concatenate all image clips into a single video
    if not image_clips:
        print("No image clips were created; please check the images directory and script.json.")
        return
    video_background = concatenate_videoclips(image_clips, method="compose")

    # Create subtitle clips based on start and end times in script.json
    subtitle_clips = []
    for entry in script_data:
        start_time = entry["start_time"]
        end_time = entry["end_time"]
        text = entry["chunk"]
        
        # Create the subtitle text clip
        subtitle_text = TextClip(
            text, fontsize=font_size, color="black", font="Arial", 
            size=(video_background.w - 40, None), method="caption"
        )

        # Create a yellow background for the subtitle text
        subtitle_bg = ColorClip(
            size=(subtitle_text.w + 20, subtitle_text.h + 10), 
            color=(255, 255, 0)  # Yellow background color
        )

        # Position the subtitle and background together at the bottom of the video
        subtitle_with_bg = CompositeVideoClip([
            subtitle_bg.set_position(("center", "bottom")),  # Yellow background
            subtitle_text.set_position(("center", "bottom"))  # Text overlay
        ], size=(video_background.w, video_background.h)).set_start(start_time).set_end(end_time)
        
        subtitle_clips.append(subtitle_with_bg)

    # Composite the final video with images, audio, and subtitles
    final_video = CompositeVideoClip([video_background, *subtitle_clips]).set_audio(audio)
    
    # Export the video
    final_video.write_videofile(output_video_path, fps=24, codec="libx264", audio_codec="aac")
    print(f"Video with images, yellow subtitle background, and audio saved to {output_video_path}")

# Usage example
output_video_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\output_with_subtitles.mp4"
generate_video_with_images_and_subtitles(final_output_path, output_script_path, output_video_path)


# Set your Pixabay API key and output folder path
# pixabay_api_key = "41610740-2e3b4e3089898192b13673058"
# video_output_folder = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\chunk_videos"

# # Ensure the video output folder exists
# if not os.path.exists(video_output_folder):
#     os.makedirs(video_output_folder)

# Function to fetch a video URL from Pixabay based on a keyword
#videos fetching and donwloading part
# def fetch_video_from_pixabay(keyword):
#     """Fetches a video URL from Pixabay based on the keyword."""
#     url = f"https://pixabay.com/api/videos/?key={pixabay_api_key}&q={keyword}"
#     response = requests.get(url)
#     if response.status_code == 200:
#         data = response.json()
#         if data["hits"]:
#             return data["hits"][0]["videos"]["medium"]["url"]  # Medium quality video URL
#     return None

# # Function to download a video from a URL
# def download_video(video_url, file_path):
    # """Downloads the video from the provided URL and saves it with the specified filename."""
    # video_response = requests.get(video_url)
    # if video_response.status_code == 200:
    #     with open(file_path, "wb") as f:
    #         f.write(video_response.content)
    #     return file_path
    # return None

# Function to trim a video to match the specified duration in script.json
#-------------------------------------------------------------> start, working hai ye use krnahai if for videos to ye use kro bs
# def trim_video(input_path, output_path, video_duration):
    # """
    # Trims the video to match the specified duration starting from 0.
    # If video_duration exceeds the video's actual duration, it trims up to the end of the video.
    # """
    # with VideoFileClip(input_path) as clip:
    #     # Use the smaller of the specified duration or the video's actual duration
    #     adjusted_end_time = min(video_duration, clip.duration)
        
    #     # Trim from the beginning up to the adjusted end time
    #     trimmed_clip = clip.subclip(0, adjusted_end_time)
    #     trimmed_clip.write_videofile(output_path, codec="libx264", audio_codec="aac")
    
    # os.remove(input_path)  # Delete the original untrimmed video file

# Function to fetch, trim, and save videos based on keywords in script.json
# def fetch_and_process_videos(script_file_path, video_output_folder):
#     """
#     Fetches, trims, and saves videos based on keywords in script.json, updating the file with video paths.
#     """
#     # Load the script JSON data
#     with open(script_file_path, 'r') as f:
#         script_data = json.load(f)

#     for chunk in script_data:
#         keyword = chunk["keyword"]
#         video_duration = chunk.get("video_duration", None)  # Get the video duration from script.json

#         if video_duration is None:
#             print(f"Skipping chunk {chunk['id']}: video_duration is missing.")
#             continue

#         # Fetch the video URL from Pixabay
#         video_url = fetch_video_from_pixabay(keyword)
        
#         if video_url:
#             # Define the file paths for the downloaded and trimmed videos
#             original_video_path = os.path.join(video_output_folder, f"chunk{chunk['id']}_original.mp4")
#             trimmed_video_path = os.path.join(video_output_folder, f"chunk{chunk['id']}_video.mp4")
            
#             # Download the video
#             if download_video(video_url, original_video_path):
#                 # Trim the video to match the specified video duration and save it
#                 trim_video(original_video_path, trimmed_video_path, video_duration)
                
#                 # Update the script JSON data with the trimmed video path
#                 chunk["video"] = trimmed_video_path
#             else:
#                 print(f"Failed to download video for keyword: {keyword}")
#         else:
#             print(f"No video found for keyword: {keyword}")

#     # Save the updated script.json with the video paths included
#     with open(script_file_path, 'w') as f:
#         json.dump(script_data, f, indent=4)
#     print("Video paths updated in script.json")

# # Run the function to fetch and process videos
# fetch_and_process_videos(output_script_path, video_output_folder)
#-------------------------------------------------------------> end

# function to concatenate all the chunk-videos into a single video
#-----------------------------------------------------------------------------------------------------> start, ye long wala hai viddeo  ka function jo chlta he but not needed
# def concatenate_clips(input_folder, output_video_path):
#     """
#     Concatenates all video clips in the input folder and saves them as a single video file.
#     """
#     # Get all video file paths from the directory
#     video_files = sorted([os.path.join(input_folder, f) for f in os.listdir(input_folder) if f.endswith('_video.mp4')])
    
#     # Load each video file as a VideoFileClip
#     clips = [VideoFileClip(file) for file in video_files]
    
#     # Concatenate the clips into one video
#     final_clip = concatenate_videoclips(clips, method="compose")
    
#     # Save the concatenated video as script.mp4
#     final_clip.write_videofile(output_video_path, codec="libx264", audio_codec="aac")
#     print(f"Concatenated video saved to {output_video_path}")

#     # Close the clip objects to free memory
#     for clip in clips:
#         clip.close()
#     final_clip.close()

# # Path to save the concatenated script video
# script_video_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\script.mp4"

# # Run the concatenation function
# concatenate_clips(video_output_folder, script_video_path)

# def generate_video_with_timed_subtitles(audio_path, script_json_path, output_video_path, background_video_path, font_size=18):
#     """
#     Generates a video with subtitles synced to the provided audio and overlays them on the background video.
#     """
#     # Load the background video
#     background_video = VideoFileClip(background_video_path)
#     video_width, video_height = background_video.size

#     # Load the audio file and get its duration
#     audio = AudioFileClip(audio_path)
#     audio_duration = audio.duration

#     # Adjust background video duration to match audio if needed
#     background_video = background_video.set_duration(audio_duration)

#     # Define subtitle styles
#     subtitle_bg_color = (255, 255, 0)  # Yellow background for subtitles
#     font_color = 'black'  # Black font color for subtitles

#     # Load subtitle timing data from script.json
#     try:
#         with open(script_json_path, "r", encoding="utf-8") as json_file:
#             script_data = json.load(json_file)
#     except FileNotFoundError:
#         print(f"Error: The file {script_json_path} does not exist.")
#         return
#     except Exception as e:
#         print(f"An error occurred while reading the JSON file: {e}")
#         return

#     # Create subtitle clips based on start and end times in script.json
#     subtitle_clips = []
#     for entry in script_data:
#         start_time = entry["start_time"]
#         end_time = entry["end_time"]
#         text = entry["chunk"]

#         # Create the subtitle text clip
#         subtitle = TextClip(text, fontsize=font_size, color=font_color, font="Arial", size=(video_width - 40, None), method="caption")

#         # Create a background for the subtitle text
#         subtitle_bg = ColorClip(size=(subtitle.w + 10, subtitle.h + 10), color=subtitle_bg_color)

#         # Overlay text on background and set timing
#         subtitle_with_bg = CompositeVideoClip(
#             [subtitle_bg.set_position(("center", "bottom")),
#             subtitle.set_position(("center", "bottom"))],
#             size=(video_width, video_height)
#         ).set_start(start_time).set_end(end_time)

#         subtitle_clips.append(subtitle_with_bg)

#     # Composite the final video with background, audio, and timed subtitles
#     final_video = CompositeVideoClip([background_video, *subtitle_clips]).set_audio(audio)
    
#     # Export the video
#     final_video.write_videofile(output_video_path, fps=24, codec="libx264", audio_codec="aac")
#     print(f"Final video with subtitles saved to {output_video_path}")

#     # Close clips to free memory
#     background_video.close()
#     audio.close()
#     for clip in subtitle_clips:
#         clip.close()
#     final_video.close()

# # Path to save the final video
# output_video_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\final-video.mp4"

# # Run the function to generate the final video with subtitles
# generate_video_with_timed_subtitles(final_output_path, output_script_path, output_video_path, script_video_path)
#-----------------------------------------------------------------------------------------------------> end

# def generate_video_with_timed_subtitles(audio_path, script_json_path, chunk_videos_folder, output_video_path, font_size=20):
#     # Load the audio file
#     audio = AudioFileClip(audio_path)
#     audio_duration = audio.duration

#     # Load subtitle timing data from script.json
#     try:
#         with open(script_json_path, "r", encoding="utf-8") as json_file:
#             script_data = json.load(json_file)
#     except FileNotFoundError:
#         print(f"Error: The file {script_json_path} does not exist.")
#         return
#     except Exception as e:
#         print(f"An error occurred while reading the JSON file: {e}")
#         return

#     # Load all chunk videos from the directory in sorted order by filename (chunk1, chunk2, etc.)
#     video_clips = []
#     for entry in script_data:
#         chunk_id = entry["id"]
#         chunk_path = os.path.join(chunk_videos_folder, f"chunk{chunk_id}_video.mp4")
        
#         if os.path.exists(chunk_path):
#             clip = VideoFileClip(chunk_path)
#             video_clips.append(clip)
#         else:
#             print(f"Warning: Video chunk file {chunk_path} not found. Skipping this chunk.")
    
#     # Concatenate all video chunks in sequence
#     if video_clips:
#         concatenated_video = concatenate_videoclips(video_clips, method="compose")
#         concatenated_video = concatenated_video.set_duration(audio_duration)  # Match duration with audio
#     else:
#         print("Error: No video chunks found to concatenate.")
#         return

#     # Create subtitle clips based on start and end times in script.json
#     subtitle_clips = []
#     for entry in script_data:
#         start_time = entry["start_time"]
#         end_time = entry["end_time"]
#         text = entry["chunk"]

#         # Create the subtitle text clip
#         subtitle = TextClip(text, fontsize=font_size, color='black', font="Arial", size=(concatenated_video.w - 40, None), method="caption")

#         # Create a background for the subtitle text
#         subtitle_bg = CompositeVideoClip([subtitle.set_position(("center", "bottom"))], size=(concatenated_video.w, concatenated_video.h)).set_start(start_time).set_end(end_time)
#         subtitle_clips.append(subtitle_bg)

#     # Combine the video, audio, and subtitles
#     final_video = CompositeVideoClip([concatenated_video, *subtitle_clips]).set_audio(audio)

#     # Export the final video
#     final_video.write_videofile(output_video_path, fps=24, codec="libx264", audio_codec="aac")
#     print(f"Final video with subtitles saved to {output_video_path}")

#     # Close clips to free resources
#     for clip in video_clips:
#         clip.close()
#     concatenated_video.close()
#     audio.close()

# output_video_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\pdf1.mp4"

# generate_video_with_timed_subtitles(final_output_path, output_script_path, video_output_folder, output_video_path)

#------------------------------------------------------------------------------------------------------> end

# final function to generate video with timed subtitles
# def generate_video_with_timed_subtitles(audio_path, script_json_path, output_video_path, font_size=18):
#     # Load the audio file
#     audio = AudioFileClip(audio_path)
#     audio_duration = audio.duration

#     # Define the video resolution and background color
#     video_width, video_height = 1280, 720
#     subtitle_bg_color = (255, 255, 0)  # Yellow background for subtitles
#     font_color = 'black'  # Black font color for subtitles

#     # Initialize video background
#     video_bg = ColorClip(size=(video_width, video_height), color=(255, 255, 255)).set_duration(audio_duration)

#     # Load subtitle timing data from script.json
#     try:
#         with open(script_json_path, "r", encoding="utf-8") as json_file:
#             script_data = json.load(json_file)
#     except FileNotFoundError:
#         print(f"Error: The file {script_json_path} does not exist.")
#         return
#     except Exception as e:
#         print(f"An error occurred while reading the JSON file: {e}")
#         return

#     # Create subtitle clips based on start and end times in script.json
#     subtitle_clips = []
#     for entry in script_data:
#         start_time = entry["start_time"]
#         end_time = entry["end_time"]
#         text = entry["chunk"]

#         # Create the subtitle text clip
#         subtitle = TextClip(text, fontsize=font_size, color=font_color, font="Arial", size=(video_width - 40, None), method="caption")

#         # Create a background for the subtitle text
#         subtitle_bg = ColorClip(size=(subtitle.w + 10, subtitle.h + 10), color=subtitle_bg_color)

#         # Overlay text on background and set timing
#         subtitle_with_bg = CompositeVideoClip([subtitle_bg.set_position(("center", "bottom")),
#                                             subtitle.set_position(("center", "bottom"))],
#                                             size=(video_width, video_height)).set_start(start_time).set_end(end_time)

#         subtitle_clips.append(subtitle_with_bg)

#     # Composite the final video with background, audio, and timed subtitles
#     final_video = CompositeVideoClip([video_bg, *subtitle_clips]).set_audio(audio)
    
#     # Export the video
#     final_video.write_videofile(output_video_path, fps=24, codec="libx264", audio_codec="aac")
#     print(f"Video with subtitles saved to {output_video_path}")

# # Usage example
# output_video_path = r"C:\Users\Happy yadav\Desktop\aura.ai\hackcbs-backend\output_with_subtitles.mp4"

# generate_video_with_timed_subtitles(final_output_path, output_script_path, output_video_path)