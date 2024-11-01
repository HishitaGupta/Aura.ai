import requests
import pdfplumber
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import os
import re
import google.generativeai as genai
from dotenv import load_dotenv
from transformers import pipeline
from moviepy.editor import VideoFileClip, ImageClip, concatenate_videoclips, CompositeVideoClip, TextClip, ColorClip, AudioFileClip
from moviepy.editor import *
from moviepy.video.fx.all import crop, loop
from moviepy.video.tools.drawing import circle
import json
import pysrt
import numpy as np
from gtts import gTTS
from googletrans import Translator
from pydub import AudioSegment
from moviepy.video.fx.all import colorx
from moviepy.video.tools.subtitles import SubtitlesClip
import base64
from docx import Document
from pptx import Presentation

# Set up folders
pictures_folder = "pictures"
videos_folder = "videos"
output_video_path = "final_slideshow.mp4"
audio_output_path = "final_audio.mp3"
audio_output_speedup_path = "final_audio_speeded_up.mp3"
srt_file_path = "subtitles.srt"

os.makedirs(pictures_folder, exist_ok=True)
os.makedirs(videos_folder, exist_ok=True)

# API URLs
img_url = 'https://pixabay.com/api/'
vid_url = 'https://pixabay.com/api/videos/'
STABILITY_API_URL = 'https://api.stability.ai/v2beta/stable-image/generate/core'
STABILITY_API_KEY = os.getenv("STABILITY_API_KEY")

# Load environment variables from .env file
load_dotenv()

API_KEY = os.getenv("IMG_API")
api_key = os.getenv("API_KEY")

# Configure API key for Google Gemini
genai.configure(api_key=api_key)
model = genai.GenerativeModel('gemini-1.5-flash')

# Set path for Tesseract OCR
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Define common resolution and frame rate
common_resolution = (1280, 720)
frame_rate = 24

def summarize_text(text):
    print("summarize_text TRIGGERED")
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    max_chunk = 1024
    text_chunks = [text[i:i + max_chunk] for i in range(0, len(text), max_chunk)]
    summary = ""
    for chunk in text_chunks:
        summary += summarizer(chunk, max_length=130, min_length=30, do_sample=False)[0]['summary_text'] + " "
    return summary.strip()

def extract_images_from_pdf(pdf_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    doc = fitz.open(pdf_path)
    image_list = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_list.append(image_bytes)
            image_path = os.path.join(output_folder, f"page_{page_num}img{img_index}.png")
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
    
    return image_list

def ocr_images_in_folder(folder_path):
    text = ""
    for filename in os.listdir(folder_path):
        if filename.endswith(".png"):
            image_path = os.path.join(folder_path, filename)
            try:
                image = Image.open(image_path)
                text += pytesseract.image_to_string(image) + "\n"
            except Exception as e:
                print(f"Error OCR processing {filename}: {e}")
    return text

def extract_text_from_pdf_images(pdf_path, output_folder):
    print("extract_text_from_pdf_images TRIGGERED")
    extract_images_from_pdf(pdf_path, output_folder)
    text = ocr_images_in_folder(output_folder)
    return text


def extract_text_from_txt(txt_path):
    with open(txt_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_pdf(pdf_path):
    print("extract_text_from_pdf TRIGGERED")
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text


def extract_text_from_file(file_path):
    if file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.txt'):
        return extract_text_from_txt(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file format")
    
def save_flow_to_json(flow_string, output_file):
    # Clean the string (not assuming it's a list)
    cleaned_string = flow_string.replace('json', '').replace('', '').strip()
    
    # Debug: Print the cleaned string to check if it's valid
    print("Cleaned string: ", cleaned_string)

    try:
        # Parse the cleaned string into a JSON object
        flow_json = json.loads(cleaned_string)
        
        # Write the JSON to file
        with open(output_file, 'w') as json_file:
            json.dump(flow_json, json_file, indent=4)
        print(f"Prompts data saved successfully to {output_file}")
    except json.JSONDecodeError as e:
        # Print the detailed error message
        print(f"Failed to decode JSON: {e}")
        print(f"Error occurred at character position: {e.pos}")
    except Exception as e:
        print(f"Failed to save JSON to file: {e}")


    
def generate_flow_from_summary(summary):
    try:
        print("generate_flow_from_summary TRIGGERED")
        inp = model.generate_content(f"""Based on the following summary, create a complete flow for a video. The flow should include the following sections:

Hooks: Generate 3 key hooks that will capture the audience's attention.
Introduction: Write 3 points introducing the policy and explaining its purpose.
Main Pointers: Identify 3 important pointers that highlight the main benefits or features of the policy.
How to Avail: Provide 3 points on how the user can avail or apply for this policy.

Format the output exactly in the following JSON structure, and surround the entire output with three backticks:
                                    
json
[
  {{
    "hooks": ["- hook point 1", "- hook point 2", "- hook point 3"],
    "intro": ["- intro point 1", "- intro point 2", "- intro point 3"],
    "pointers": ["- main pointer 1", "- main pointer 2", "- main pointer 3"],
    "howToAvail": ["- how to avail 1", "- how to avail 2", "- how to avail 3"]
  }}
]

do not add asterisks in this, only add hyphen before the points.
                                    
Summary: {summary}""")
        flow = inp.text
        print("gemini flow: ", flow)
        
        # Pass the flow (string) to save_flow_to_json
        save_flow_to_json(flow, r"D:\hackerx\Phosphenes-HackRx-5.0\hackrx-backend\flow.json")
    except Exception as e:
        print(f"Error generating speech: {e}")

generate_flow_from_summary("Bajaj Allianzs Health Infinity Policy comes with new comprehensive benefits at competitive premiums. Indemnity Health Insurance Plan without any limit on Sum Insured. Road Ambulance up to Rs. 5000 per hospitalization. Preventive Health Check Up at the end of every 3 policy years. No prepolicy medical tests up to 45 years of age subject to clean proposal form Preexisting disease covered after 36 months from your first Health Policy. Income tax benefit under 80 D of the IT Act on premiums paid for this policy. Room Rent options under the policy. Policy can be taken for 1year 2years OR 3years. Policy can be opted by NonResident Indians including PIOs Persons of Indian Origin and OCIs Overseas Citizens of India Policy will be issued during their stay in India and premium is paid in Indian currency. The Company will pay the Insured,Reasonable and Customary Medical Expenses incurred subject to certain conditions. If admitted in ICU, the company will pay up to ICU actual expenses provided by Hospital. The Medical Expenses incurred during the 60 days immediately before the Insured was Hospitalized, provided that Such Medical Exp expenses were incurred for the same illnessinjury for which subsequent Hospitalization wasrequired. Relevant laboratory diagnostic tests, Xray and such similar expenses that are medically prescribed by the treat tumultuous Medical Practitioner. The Company will pay the reasonable cost upto a maximum of Rs 5000 per Hospitalization incurred on an ambulance. Claim under this section shall be payable by the Company only when. Such life threatening emergency condition is certified by the Medical Practitioner, and. The Company has accepted Insureds Claim under Inpatient Hospitalization Treatment Point No. 1 or DayCare Procedures section Point. 5 of the Policy. The Company will pay the Insured medical expenses as listed above under Inpatient Hospitalization Treat ment Point No. 1 for Day care procedures  Surgeries taken as an inpatient in a hospital or day care centre but not in the outpatient department. The Company will not be liable to make any payment for any claim directly or indirectly caused by, based on, or attributable to any of the following. Expenses related to the treatment of a preexisting Disease PED and its direct complications shall be excluded until the expiry of 36 months of continuous coverage. Expenses related to the treatment of the listed Conditions, surgeriestreatments shall be excluded until theexpiry of 24 months of continuous coverage after the date of inception of the first Health Infinity Policy with Us. After 24 months of continuous coverage, the maximum limit for each claim will be restricted to 100 times per day. List of specific diseasesprocedures is as below. The limit of indemnity will be applicable only for the procedures. The limit of indemnity for each claim will be restricted to 100 times per day room rent limit for the below listed conditions. Expenses related to the treatment of any illness within 30 days from the first Policy commencement date shall be excluded except claims arising due to an accident. Any dental treatment that comprises of cosmetic. surgery, dentures, dental prosthesis, dental implants, ortho.dontics, surgery of any kind unless as a. result of Accidental Bodily Injury to natural teeth and also requiring. hospitalization. War, invasion, acts of foreign enemies, hostilities whether war be declared or not, civil war, commotion, unrest, unrest. Any Medical expenses incurred due to Act of Terrorism will be covered under the Policy. Circumcision unless required for the treatment of Illness or Accidental bodily injury, is excluded. Expenses related to the surgical treatment of obesity that does not fulfil all the below conditions. Surgery to be conducted is upon the advice of the Doctor. The surgeryProcedure conducted should be supported by clinical protocols. The cost of spectacles, contact lenses, hearing aids, crutches, dentures, artificial teeth and all other external medical devices. Expenses related to any treatment necessitated due to participation as a professional in hazardous or adventure sports. Expenses for treatment directly arising from or consequent upon any Insured committing or attempting to commit a breach of law with criminal intent. Equipment of any kind used at home as post Hospitalization care. expenses up to the stage of stabilization are payable but not the complete claim. Treatments received in heath hydros, nature cure clinics, spas or similar establishments. Expenses related to the treatment for correction of eye sight due to refractive error less than 7.5 dioptres. Unproven treatments are treatments, procedures or supplies that lack significant medical documentation to support their effectiveness. Medical tests are applicable for members 46 years and above as per grid given below. Treatment for any other system other than modern medicine also known as Allopathyand AYUSH therapies is not covered. 46 years and above Room rent from 25,000 Medical Tests required as listed below FMR, CBC, Urine R, ECG, Lipid profile, Fasting BSL, HBA1c,SGOT, SGPT, Sr.Creatinine, Total Protein, Jr. Globulin, AG Ratio, USG abdomen  Pelvis, TMT. 20 discount on published premium rates to employees of Bajaj Allianz  its group companies. 5 discount is extended for the policies purchased online through websitedirect customers. The waiting periods specified in the Section CI.1 to I.3 shall be reduced by the number of continuous preceding years of coverage of the Insured under the previous health insurance Policy. Cost Sharing applicable if admission in higher room category. The policy covers all hospitalization expenses incurred during the policy period subject to the policy coverage, conditions, definitions  exclusions. The maximum limit of indemnity for ailmentsconditions as mentioned in waiting period A3  A4 will be rericted to 100 times the room rent limit opted for each claim. Copayment will be applicable for claims for both Network and Non Network Hospitals. Copayment as defined in point i would apply first followed by that defined under point ii g 100 times the room rent limit would be not be admissible under the policy. He has opted for a private room with per day room rent of 4000 INR, which is within the room rent limit opted. As per the policy condition there is no copayment up to 100 times of room rent. There is noCopay up to 500,000 INR. Policy copayment applicable on 250,000 INR ij and 20 on the final admissible amount exceeding 500,000. Room rent in which Mr Amit is hospitalized  6,000 InR per day. Total Claimed Amou 692,000 Nt  750,000 INR. Break up of Hospital Bill Room rent 600010 days 60,000. Professional Charges Doctors, Surgeons etc. 300,000c Investigations 100,000d OT, nursing etc. charges100,000e Medicines and consumables 100,00f Internal implant cost 80,000g Other NonMedical Expenses Not Payable 10,000h Total Bill amount 750,00,000 in India. Policy copayment applicable on 146,648 o  p 29,333r Total amount payable by insurer m  q 617,333 per day. If you have not made any claim during the Free look period, you shall be entitled to refund of premium subject to the terms and conditions of this Policy. You have the option of canceling the Policy stating the reasons for cancellation. The Free Look period is not applicable for renewal policies. Claims under other policies may be made after exhaustion of Sum Insured in the earlier chosen Policy. If the amount to be claimed exceeds the sum under a single policy, the policyholder has the right to choose which policy they want to claim from. Under normal circumstances, renewal will not be refused except on the grounds of Your moral hazard, misrepresentation, fraud, or your noncooperation. If Insured has multiple Policies, he she has the right to prefer claims from other PolicyPolicies. A grace period of 30 days is permissible. If the insured person has opted for Payment of Premium on an instalment basis i.e. Annual, Half Yearly, Quarterly or Monthly, the following Conditions shall apply notwithstanding any terms contrary elsewhere in the policy The grace period of fifteen days where premium is paid on a monthly instalments is available on the premium due date. We will not apply any additional loading on your policy premium at renewal based on claim experience. CIN U66010PN2000PLC015329, UIN BAJHLIP21005V022021 13. The insured person will get the accrued continuity benefit in respect of the Waiting Periods, Specific Waiting Periods. in the event of payment of premium within the stipulated grace Period. Cancellation of policy where full premium received at policy inception Annual Policy The premium refund for the unexpired risk period will be on a prorata basis, provided no claim has been made during the policy year. The Company may cancel the Policy at any time on the grounds of misrepresentation, nondisclosure of materialfacts, or fraud by the Policyholder. There will be no refund of premiums for cancellations on these grounds. Portability shall be allowed under all individual indemnity health insurance policies issued by General Insurers. Migration of health insurance policies is an option for all policy holders. Policyholders can choose to leave their individual or group health insurance policy at the time of exit. Migration from group policies to individual policy will be subject to underwriting. There is possibility of withdraw of the group policy. Migration shall be applicable to the extent of the sum insured under the previous policy. You can choose among our available similar and closely similar Health insurance products. Upon Your so choosing Our new product, You will be charged the Premium as per Our Underwriting Policy for such chosen new product. We reserve Our right to do so with a intimation of 3 months to all the existing insured members. All Claims will be settled by In house claims settlement team of the company and no TPA is engaged. Cashless treatment is only available at Network Hospitals. The Companys representative then within 2 hours will respond with Approval, Rejection or an more information. If the procedure above is followed, the Insured will not be required to directly pay for the bill amount in the horriblyNetwork Hospital that the Company is liable under Section A1 InPatient Hospitalization Treatment above. Preauthorization does not guarantee that all costs and expenses will be covered. Original bills and evidence of treatment in respect of the same shall be left with the Network Hospital. The Insured or someone claiming on his her behalf must promptly and in any event within 30 days of dischargefrom a Hospital give the Company documentation. The Insured must have himself  herself examined by the Companys medical advisors if the Company ask for it. eInsured may provide the Company with the attested Xerox copies of such documents along with a declaration from the particular insurer specifying the availability of the original copies. Waiver of conditions i and vi may be considered in extreme cases of hardship. Medical Practitioner certificate to transfer the Injured person to a higher medical centre for further treatment. Cashless settlement letter or other company settlement letter. First consultation letter for the current ailment. Original bills towards Investigations done  Laboratory Bills. We will make payment when You or someone claiming on Your behalf has provided Us with necessary documentation and information. If there is no Nominee and You are incapacitated or deceased, we will pay Your heir, executor or validly appointed legal representative. In the cases of delay in the payment, the Company shall be liable to pay interest at a rate which is 2 above the bank rate prevalent at the beginning of the financial year. The Company will settle the claim within 45 days from the date of receipt of last necessary document. In case of delay beyond stipulated 45 days, the Company will be liable to pay interest at a rate of 2 above the bank rate. Cashless facility offered through network hospitals of Bajaj Allianz only. Policy Schedule, Policy Wordings, Cashless Cards and Health Guide will be sent to your mailing addressmentioned on the proposal form. Senior Citizen Cell for Insured Person who are over 60 years of age. For more details on the coverage, terms andclusions, please get in touch with nearest office of Bajaj Allianz. Policy holders can download Insurance Wallet for one touch access. To update your contact details i.e., Mobile No., Email ID, PAN Card, and Bank Account details, please use chatbot, visit our website.")

def fetch_flow(file_path):
    """Fetches the flow JSON structure from a file."""
    with open(file_path, 'r') as f:
        return json.load(f)

def text_to_speech(text, output_file):
    """Convert text to speech using gTTS and save as an audio file."""
    tts = gTTS(text)
    tts.save(output_file)
    return output_file

# def generate_text_clip(text, duration, font_size=40, color='white', bg_color='black', font='Arial'):
#     """Generate a text clip with MoviePy."""
#     text_clip = TextClip(text, fontsize=font_size, color=color, font=font, bg_color=bg_color)
#     text_clip = text_clip.set_duration(duration).set_position('center')
#     return text_clip

# def create_section_clip(text_list, section_audio, section_image=r"D:\hackerx\Phosphenes-HackRx-5.0\hackrx-backend\ai_generated_images\bg_img.png", duration=5):
#     """Creates a clip for a section with optional background image and text overlay."""
#     clips = []
    
#     # Load the section's audio (voiceover)
#     audio_clip = AudioFileClip(section_audio)
    
#     # Create video clips for each text item
#     for text in text_list:
#         text_clip = generate_text_clip(text, duration)
#         if section_image:
#             # Combine with background image if provided
#             image_clip = ImageClip(section_image).set_duration(duration)
#             combined_clip = CompositeVideoClip([image_clip, text_clip])
#         else:
#             combined_clip = text_clip
#         combined_clip = combined_clip.set_audio(audio_clip)
#         clips.append(combined_clip)
    
#     # Concatenate all text clips into one section video
#     section_video = concatenate_videoclips(clips)
#     return section_video


# def generate_video_from_flow(flow_json, output_video_path, background_music=None):
#     """Generates a video based on the JSON flow."""
#     clips = []
    
#     # Iterate through the flow sections
#     for section in flow_json:
#         # Create hooks section
#         hooks = section.get('hooks', [])
#         hooks_text = ' '.join(hooks)
#         hooks_audio = text_to_speech(hooks_text, 'hooks.mp3')
#         hooks_clip = create_section_clip(hooks, hooks_audio, duration=3)
#         clips.append(hooks_clip)
        
#         # Create intro section
#         intro = section.get('intro', [])
#         intro_text = ' '.join(intro)
#         intro_audio = text_to_speech(intro_text, 'intro.mp3')
#         intro_clip = create_section_clip(intro, intro_audio, duration=4)
#         clips.append(intro_clip)
        
#         # Create main pointers section
#         pointers = section.get('pointers', [])
#         pointers_text = ' '.join(pointers)
#         pointers_audio = text_to_speech(pointers_text, 'pointers.mp3')
#         pointers_clip = create_section_clip(pointers, pointers_audio, duration=5)
#         clips.append(pointers_clip)
        
#         # Create how to avail section
#         how_to_avail = section.get('howToAvail', [])
#         avail_text = ' '.join(how_to_avail)
#         avail_audio = text_to_speech(avail_text, 'how_to_avail.mp3')
#         avail_clip = create_section_clip(how_to_avail, avail_audio, duration=4)
#         clips.append(avail_clip)
    
#     # Concatenate all section clips
#     final_video = concatenate_videoclips(clips)
    
#     # Add background music (if provided)
#     if background_music:
#         music_clip = AudioFileClip(background_music)
#         final_video = final_video.set_audio(CompositeAudioClip([final_video.audio, music_clip]))
    
#     # Write the final video file
#     final_video.write_videofile(output_video_path, fps=24)


# flow_data = fetch_flow(r'D:\hackerx\Phosphenes-HackRx-5.0\hackrx-backend\flow.json')
# generate_video_from_flow(flow_data, 'final_video.mp4', background_music=r'D:\hackerx\Phosphenes-HackRx-5.0\hackrx-backend\background.mp3')
import numpy as np
from moviepy.editor import *
from math import sin, pi

def generate_final_video_with_bouncing_images(flow_data, background_audio_path, output_file="final_video.mp4"):
    width, height = 1920, 1080
    font_size = 50
    heading_font_size = 80
    fade_duration = 0.5
    duration_per_point = 3
    section_pause = 1  # Pause between sections
    top_margin = 100
    left_margin = 50
    spacing = 100

    # Set max width for text to cover 70% of the screen on the left
    text_max_width = int(width * 0.7) - left_margin

    # Load the background image
    background = ImageClip(r"C:\Users\Happy yadav\Desktop\aura.ai\hackrx-backend\ai_generated_images\bg_img.png").resize((width, height))

    # Load the illustration (can be an image, gif, or video) and resize it
    illustration = ImageClip(r"C:\Users\Happy yadav\Desktop\aura.ai\hackrx-backend\ai_generated_images\bg_img.png").resize(height=500)

    # Apply a rounded corner effect to the illustration (circle effect)
    radius = min(illustration.size) // 2  # Half the smallest dimension for circular effect
    illustration = illustration.on_color(color=(0, 0, 0), col_opacity=0, size=(illustration.w + 2*radius, illustration.h + 2*radius))

    # Apply the circular border using moviepy's crop method
    illustration = illustration.crop(x1=radius, y1=radius, x2=illustration.w - radius, y2=illustration.h - radius)

    # Set bounce parameters (amplitude and frequency)
    bounce_amplitude = 20  # How high the image bounces
    bounce_frequency = 2  # How fast the bounce happens (higher = faster)

    # Create a lambda function to apply the bouncing animation
    def bounce_position(t):
        x_pos = width - illustration.w - 50  # Fixed position on the right side
        y_pos = height // 2 + bounce_amplitude * sin(bounce_frequency * t * 2 * pi)  # Bouncing effect
        return x_pos, y_pos

    # Apply the bounce effect to the illustration (moving it every frame)
    illustration = illustration.set_position(lambda t: bounce_position(t))

    # List to hold all section clips
    section_clips = []

    # Dictionary for section headings
    section_headings = {
        "intro": "About",
        "pointers": "Benefits",
        "howToAvail": "How to Avail"
    }

    for section in flow_data:
        for section_name, section_points in section.items():
            section_text_clips = []
            current_y = top_margin

            if section_name in section_headings:
                heading_text = section_headings[section_name]

                # Create the heading clip
                heading_clip = TextClip(
                    heading_text, fontsize=heading_font_size, color='white', font='Arial-Bold', kerning=5
                )
                heading_clip = heading_clip.set_position(('center', top_margin - 50))  # Top center
                heading_clip = heading_clip.set_start(0).set_duration(len(section_points) * duration_per_point + section_pause)
                heading_clip = heading_clip.fadein(fade_duration).fadeout(fade_duration)

                section_text_clips.append(heading_clip)

                current_y += heading_clip.h + 50

            for i, point in enumerate(section_points):
                # Create the text clip, wrapping and setting max width for 70% left space
                txt_clip = TextClip(
                    point, fontsize=font_size, color='white', font='Arial-Bold', kerning=5,
                    method='caption', size=(text_max_width, None)  # Limit text width to 70% of the screen
                )
                txt_clip = txt_clip.set_position((left_margin, current_y))  # Position on the left part

                start_time = i * duration_per_point
                txt_clip = txt_clip.set_start(start_time).set_duration((len(section_points) - i) * duration_per_point + section_pause)

                txt_clip = txt_clip.fadein(fade_duration)

                section_text_clips.append(txt_clip)
                current_y += txt_clip.h + spacing

            # Create a composite clip for the section
            section_duration = len(section_points) * duration_per_point + section_pause
            section_clip = CompositeVideoClip(section_text_clips, size=(width, height))
            section_clip = section_clip.set_duration(section_duration)

            section_clip = section_clip.fadeout(fade_duration)

            section_clips.append(section_clip)

    # Concatenate all section clips
    final_video = concatenate_videoclips(section_clips)

    # Combine the background, final video, and illustration (illustration will be over everything)
    video = CompositeVideoClip([background.set_duration(final_video.duration), final_video, illustration.set_duration(final_video.duration)])

    # Load and lower the volume of the background audio
    background_audio = AudioFileClip(background_audio_path).volumex(0.2)
    background_audio = background_audio.set_duration(video.duration)

    # Set the final audio on the video clip
    video = video.set_audio(background_audio)

    # Set the fps value
    video.fps = 24

    # Write the video to a file
    video.write_videofile(output_file, codec='libx264', audio_codec='aac', fps=24)

flow_data = fetch_flow(r'C:\Users\Happy yadav\Desktop\aura.ai\hackrx-backend\flow.json')
generate_final_video_with_bouncing_images(flow_data, r"C:\Users\Happy yadav\Desktop\aura.ai\hackrx-backend\background.mp3")





def extract_chunks_from_srt(srt_file_path):
    subs = pysrt.open(srt_file_path)
    chunks = [sub.text for sub in subs]
    return chunks

def generate_keywords_from_summary(summary, srt_file_path):
    speeches = ""
    keywords = []

    # Generate speech from summary
    try:
        inp = model.generate_content(f"Write a minimum exact 250 and maximum exact 300 words more accurate summary based on the previous summary in paragraph format and it should be plain text with no bullet points, no '/n' and no bold stuff, i am using it as input for my tts: {summary}")
        speeches = inp.text
    except Exception as e:
        print(f"Error generating speech: {e}")

    # Extract chunks from SRT and generate keywords
    try:
        chunks = extract_chunks_from_srt(srt_file_path)

        for chunk in chunks:
            try:
                # Generate keywords for each chunk
                res = model.generate_content(f"Generate 1 unique, main and relevant keyword make sure it's one word and relevant enough to generate the video  based on this 1 chunk: {chunk}")
                keywords += re.findall(r'\\(.?)\\*', res.text)
            except Exception as e:
                print(f"Error generating keywords for chunk: {chunk}, Error: {e}")

        # Deduplicate keywords
        keywords = list(set(keywords))
    except Exception as e:
        print(f"Error extracting keywords from SRT: {e}")

    return {
        "speech": speeches,
        "keywords": keywords,
    }


# def generate_keywords_from_summary(summary):
#     speeches = ""
#     keywords = []
    

#     try:
#         inp = model.generate_content(f"Write a minimum exact 250 and maximum exact 300 words more accurate summary based on the previous summary in paragraph format and it should be plain text with no bullet points, no '/n' and no bold stuff, i am using it as input for my tts: {summary}")
#         speeches = inp.text
#     except Exception as e:
#         print(f"Error generating speech: {e}")



#     try:
#         res = model.generate_content(f"Generate 10 unique, main and relevant keywords based on summary make sure it's one word and relevant enough to generate an image which I can use in making video: {summary}")
#         keywords = re.findall(r'\\(.?)\\*', res.text)
#     except Exception as e:
#         print(f"Error generating keywords: {e}")

#     return {
#         "speech": speeches,
#         "keywords": keywords,
        
#     }

# Subtitle part

def generate_subtitles_from_speech(speech_text, audio_duration, output_srt_path, chunk_duration=5):
    # Split speech into chunks
    words = speech_text.split()
    chunk_size = int(audio_duration // chunk_duration)
    chunks = [words[i:i + chunk_size] for i in range(0, len(words), chunk_size)]
    
    # Calculate total number of chunks
    total_chunks = len(chunks)
    
    # Adjust chunk_duration if it exceeds the audio duration
    if total_chunks > 1 and (audio_duration / total_chunks) < chunk_duration:
        chunk_duration = audio_duration / total_chunks
    
    # Create subtitles
    subs = pysrt.SubRipFile()
    for i, chunk in enumerate(chunks):
        start_time = i * (audio_duration / total_chunks)
        end_time = (i + 1) * (audio_duration / total_chunks)
        
        # Ensure end time does not exceed audio duration
        if end_time > audio_duration:
            end_time = audio_duration
        
        start_time = pysrt.SubRipTime(seconds=start_time)
        end_time = pysrt.SubRipTime(seconds=end_time)

        subtitle = pysrt.SubRipItem(index=i+1, start=start_time, end=end_time, text=' '.join(chunk))
        subs.append(subtitle)

    # Save the subtitles to a file
    subs.save(output_srt_path, encoding='utf-8')
    print(f"Subtitles saved to {output_srt_path}")

def generate_prompts_from_srt_chunks(srt_file_path):
    chunks = extract_chunks_from_srt(srt_file_path)
    prompts = []

    for chunk in chunks:
        try:
            # Generate a prompt for each chunk using the model
            inp = model.generate_content(f"Generate 1 Unique and realistic prompt containing 'description' in JSON list fromat for generating the images using AI based on the this 1 chunk: {chunk}")
            prompt_text = inp.text
            prompts.append(prompt_text)
        except Exception as e:
            print(f"Error generating prompt for chunk: {chunk}, Error: {e}")
    
    return prompts

def create_dynamic_pointers_from_srt(srt_file_path, image_path, output_video_path, duration=5, vertical_margin=50):
    # Load the SRT file and extract chunks directly in one step
    subs = pysrt.open(srt_file_path)
    chunks = [sub.text for sub in subs]

    # Generate pointers by grouping every 3 chunks into one set of 3 chunks
    pointer_texts = []
    for i in range(0, len(chunks), 3):
        # Join every 3 chunks into one string
        chunk_group = " ".join(chunks[i:i+3])
        print(chunk_group)
        # Generate a pointer from Gemini AI based on the chunk group
        try:
            prompt = f"Generate a concise 10-word pointer summarizing the following text: {chunk_group}"
            response = model.generate_content(prompt)
            pointer_text = response.text.strip()  # Get the generated pointer text
            pointer_texts.append(pointer_text)  # Add it to the pointer list
        except Exception as e:
            print(f"Error generating pointer with Gemini AI: {e}")
            pointer_texts.append("Error generating pointer")  # If error, append a placeholder text

    # Load the image
    image_clip = ImageClip(image_path, duration=duration)

    # Create a list to hold the text clips
    text_clips = []

    # Get image height to position the text pointers vertically
    image_width, image_height = image_clip.size

    # Starting y-position for the first pointer
    y_start = image_height // 4

    # Add each text as a clip and position them vertically starting from the left
    for idx, text in enumerate(pointer_texts):
        # Create text clip
        text_clip = TextClip(text, fontsize=35, color='white', font="Arial-Bold")

        # Calculate the position: Left-aligned and vertically stacked
        x_pos = 50  # Starting from the left with a 50px margin
        y_pos = y_start + idx * (text_clip.size[1] + vertical_margin)  # Stacked vertically with margin

        # Set the position and duration for the text clip
        text_clip = text_clip.set_position((x_pos, y_pos)).set_duration(duration)

        # Apply the fade-in effect with the text fully invisible at the start
        text_clip = text_clip.crossfadein(1.5)  # Smooth fade-in over 1.5 seconds

        # Add a delay for each text clip to make it fade in at different times
        text_clip = text_clip.set_start(idx * 0.5)  # Delay between fades

        text_clips.append(text_clip)

    # Combine the image and text clips
    final_clip = CompositeVideoClip([image_clip] + text_clips).set_duration(duration)

    # Save the video
    final_clip.write_videofile(output_video_path, fps=24)



#gen prompts from summmary
# def generate_prompts_from_summary(summary):
#     promp_string = ""
#     try:
#         inp = model.generate_content(f"Generate 10 Unique and realistic prompt containing 'description' in JSON list fromat for generating the images using AI based on the this summary: {summary} ")
#         promp_string = inp.text.strip("\n")
#         print(promp_string)
#     except Exception as e:
#         print(f"Error generating prompts: {e}")
#     return promp_string


# def save_prompts_to_json(promp_string, output_file):
#     # cleaned_string = promp_string.replace('json', '').replace('', '').strip()
#     cleaned_string = [s.replace('json', '').replace('', '').strip() for s in promp_string]
    
#     try:
#         prompts_json = json.loads(cleaned_string)
#         with open(output_file, 'w') as json_file:
#             json.dump(prompts_json, json_file, indent=4)
#         print(f"Prompts data saved successfully to {output_file}")
#     except json.JSONDecodeError as e:
#         print(f"Failed to decode JSON: {e}")

def save_prompts_to_json(promp_string, output_file):
    # Clean each string in the list
    cleaned_strings = [s.replace('json', '').replace('', '').strip() for s in promp_string]

    # Initialize an empty list to store the parsed JSON objects
    prompts_json_list = []

    # Process each cleaned string
    for cleaned_string in cleaned_strings:
        try:
            # Parse the cleaned string into a JSON object
            prompts_json = json.loads(cleaned_string)
            prompts_json_list.append(prompts_json)
        except json.JSONDecodeError as e:
            print(f"Failed to decode JSON in string: {cleaned_string}\nError: {e}")
    
    # Save the list of JSON objects to the output file
    try:
        with open(output_file, 'w') as json_file:
            json.dump(prompts_json_list, json_file, indent=4)
        print(f"Prompts data saved successfully to {output_file}")
    except Exception as e:
        print(f"Failed to save JSON to file: {e}")
# Quiz part

def generate_quiz(text):
    quiz_string = ""
    try:
        inp = model.generate_content(f"Generate quiz which contains 3 questions with unique answers in MCQ format containing 'question', 'options', 'answer' in JSON list format based on the text: {text}")
        print(inp.text)
        quiz_string = inp.text
    except Exception as e:
        print(f"Error generating quiz: {e}")
    return quiz_string

def save_quiz_to_json(quiz_string, output_file):
    cleaned_string = quiz_string.replace('json', '').replace('', '').strip()
    
    try:
        quiz_json = json.loads(cleaned_string)
        with open(output_file, 'w') as json_file:
            json.dump(quiz_json, json_file, indent=4)
        print(f"Quiz data saved successfully to {output_file}")
    except json.JSONDecodeError as e:
        print(f"Failed to decode JSON: {e}")

def clean_text(text):
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^a-zA-Z0-9\s.,]', '', text)
    return text.strip()

def get_text_from_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as text_file:
            text = text_file.read()
        return text
    except FileNotFoundError:
        print(f"The file {file_path} does not exist.")
        return None
    except Exception as e:
        print(f"An error occurred while reading the file: {e}")
        return None

def ask_aura_question(user_question, txt_file_path):
    answer = ""
    try:
        # Get the extracted text from the txt file
        extracted_text = get_text_from_txt(txt_file_path)
        
        if not extracted_text:
            return "Error: Could not retrieve text from file."

        # Use the extracted text and the user's question to prompt Gemini AI
        prompt = f"Answer the following question based on this text: '{extracted_text}'.\nQuestion: {user_question}"
        response = model.generate_content(prompt)
        print("Answer response: ", response.text)  # Debug print for checking the response
        answer = response.text
    except Exception as e:
        print(f"Error while asking Gemini AI: {e}")
        return f"Error: {e}"
    
    return answer


# def save_image_from_url(image_url, save_directory, image_index):
#     os.makedirs(save_directory, exist_ok=True)
    
#     try:
#         image_response = requests.get(image_url)
#         if image_response.status_code == 200:
#             image_path = os.path.join(save_directory, f'image_{image_index}.jpg')
#             with open(image_path, 'wb') as file:
#                 file.write(image_response.content)
#             print(f"Image {image_index} saved as {image_path}")
#         else:
#             print(f"Failed to download image {image_index}. Status code: {image_response.status_code}")
#     except Exception as e:
#         print(f"Error saving image {image_index}: {e}")

def save_video_from_url(video_url, save_directory, video_index):
    os.makedirs(save_directory, exist_ok=True)
    
    try:
        video_response = requests.get(video_url)
        if video_response.status_code == 200:
            video_path = os.path.join(save_directory, f'video_{video_index}.mp4')
            with open(video_path, 'wb') as file:
                file.write(video_response.content)
            print(f"Video {video_index} saved as {video_path}")
        else:
            print(f"Failed to download video {video_index}. Status code: {video_response.status_code}")
    except Exception as e:
        print(f"Error saving video {video_index}: {e}")

def trim_video(video_path, duration=7):
    try:
        with VideoFileClip(video_path) as video:
            print(f"Original Duration: {video.duration}")
            if video.duration > duration:
                trimmed_video = video.subclip(0, duration)
                trimmed_video_path = video_path.replace(".mp4", "_trimmed.mp4")
                trimmed_video.write_videofile(trimmed_video_path, codec='libx264', audio_codec='aac')
                print(f"Trimmed Duration: {trimmed_video.duration}")
                os.replace(trimmed_video_path, video_path)
    except Exception as e:
        print(f"Error trimming video {video_path}: {e}")
        if os.path.exists(video_path):
            os.remove(video_path)

def generate_image_from_prompt(prompt, images_folder, image_index):
    headers = {
        'Authorization': f'Bearer sk-Xips2HbQGXa2GTbRE3Y08s933bqp7eNhSUGw4eyFW1Vu4Cep',
        
        'accept': 'image/*'
    }

    data = {
        "prompt": prompt,
        # "aspect_ratio": "16:9",
        "output_format": "png",
        "aspect_ratio": "16:9",
    }
    print("data: ", data['prompt'])

    files = {
        'prompt': (None, prompt)  # Use None to indicate that it's not a file
    }

    print("prompt: ", files['prompt'])
    try:
        print("Sending prompt:", prompt)
        response = requests.post(
            "https://api.stability.ai/v2beta/stable-image/generate/core",
            headers=headers,
            files=files,  # Sending as multipart/form-data
            data=data
        )

        # Check the response status
        print("Response Status Code:", response.status_code)
        # print("Response Content:", response.text)

        response.raise_for_status()  # Raise an error for bad responses

        # If receiving the image directly
        if headers['accept'] == 'image/*':
            image_data = response.content
            image_path = os.path.join(images_folder, f"generated_image_{image_index}.png")
            with open(image_path, "wb") as img_file:
                img_file.write(image_data)
            print(f"Image saved: {image_path}")
        else:
            # If receiving base64 JSON response
            json_response = response.json()
            base64_image = json_response.get('image', None)
            if base64_image:
                image_data = base64.b64decode(base64_image)
                image_path = os.path.join(images_folder, f"generated_image_{image_index}.png")
                with open(image_path, "wb") as img_file:
                    img_file.write(image_data)
                print(f"Image saved: {image_path}")
            else:
                print("No image data received in JSON response.")

    except requests.exceptions.HTTPError as http_err:
        print(f"HTTP error occurred: {http_err}")
    except Exception as e:
        print(f"Error generating image from prompt: {e}")




#pixaby code 


def generate_and_save_images_and_videos_for_keywords(keywords):
    # Ensure the headers are set up correctly
    headers = {'Authorization': f'Bearer {API_KEY}'}
    
    # Loop through keywords to fetch images and videos
    for i, keyword in enumerate(keywords):
        print(f"Processing keyword {i + 1}: {keyword}")

        params = {
            'key': API_KEY,
            'q': keyword,
            'video_type': 'animation',    # For videos
            'per_page': 3             # Number of results per page
        }

        try:
            if i < 10:  # Fetch videos
                response = requests.get(vid_url, params=params)
                if response.status_code == 200:
                    data = response.json()
                    videos = data.get('hits', [])
                    for idx, video in enumerate(videos):
                        video_url = video.get('videos', {}).get('medium', {}).get('url', '')
                        if video_url:
                            save_video_from_url(video_url, videos_folder, i + 1 + idx)
                            video_path = os.path.join(videos_folder, f'video_{i + 1 + idx}.mp4')
                            trim_video(video_path)
                else:
                    print(f"Failed to fetch videos for keyword {i + 1}. Status code: {response.status_code}")

            else:
                print("Nothing is fetching...")
            # else:  # Fetch videos
            #     response = requests.get(vid_url, params=params)
            #     if response.status_code == 200:
            #         data = response.json()
            #         videos = data.get('hits', [])
            #         for idx, video in enumerate(videos):
            #             video_url = video.get('videos', {}).get('medium', {}).get('url', '')
            #             if video_url:
            #                 save_video_from_url(video_url, videos_folder, i + 1 + idx)
            #                 video_path = os.path.join(videos_folder, f'video_{i + 1 + idx}.mp4')
            #                 trim_video(video_path)
            #     else:
            #         print(f"Failed to fetch videos for keyword {i + 1}. Status code: {response.status_code}")
        except Exception as e:
            print(f"Error processing keyword {keyword}: {e}")

def clean_up_videos(folder_path, max_duration=7):
    for filename in os.listdir(folder_path):
        if filename.endswith(".mp4"):
            video_path = os.path.join(folder_path, filename)
            try:
                with VideoFileClip(video_path) as video:
                    if video.duration > max_duration:
                        os.remove(video_path)
                        print(f"Deleted {video_path} because its duration was greater than {max_duration} seconds.")
            except Exception as e:
                print(f"Error checking video duration for {video_path}: {e}")



def create_slideshow_with_audio(images_folder, videos_folder, output_video_path, audio_path, overlay_video_path, srt_file_path, image_duration=2, fade_duration=1):
    image_clips=[]
    video_clips = []

    # Calculate the total duration of the audio
    audio_clip = AudioFileClip(audio_path)
    audio_duration = audio_clip.duration
    audio_clip.close()

    # Calculate the number of images and videos
    num_images = len([f for f in os.listdir(images_folder) if f.endswith(('.jpg', '.png'))])
    num_videos = len([f for f in os.listdir(videos_folder) if f.endswith('.mp4')])

    # Calculate the duration per image and video to match the audio length
    if  num_images + num_videos > 0:
        duration_per_clip = audio_duration / (num_images+num_videos)

    # Load and process image clips
    for filename in sorted(os.listdir(images_folder)):
        if filename.endswith(('.jpg', '.png')):
            image_path = os.path.join(images_folder, filename)
            image_clip = ImageClip(image_path, duration=duration_per_clip).set_fps(frame_rate)
            
            # Resize while maintaining aspect ratio
            image_clip = image_clip.resize(height=common_resolution[1])
            image_clip = image_clip.set_duration(duration_per_clip).fadein(fade_duration).fadeout(fade_duration)

            background = ColorClip(size=common_resolution, color=(0, 0, 0), duration=duration_per_clip)
            image_clip = CompositeVideoClip([background, image_clip.set_position("center")])
            image_clips.append(image_clip)

    # Load and process video clips
    for filename in sorted(os.listdir(videos_folder)):
        if filename.endswith('.mp4'):
            video_path = os.path.join(videos_folder, filename)
            video_clip = VideoFileClip(video_path).set_fps(frame_rate)

            video_clip = video_clip.resize(height=common_resolution[1])
            video_clip = video_clip.set_duration(duration_per_clip).fadein(fade_duration).fadeout(fade_duration)

            background = ColorClip(size=common_resolution, color=(0, 0, 0), duration=duration_per_clip)
            video_clip = CompositeVideoClip([background, video_clip.set_position("center")])
            video_clips.append(video_clip)

    # Create the final list of clips by alternating between images and videos
    clips = []
    max_len = max(len(image_clips), len(video_clips))
    # max_len = len(video_clips)

    for i in range(max_len):
        if i < len(image_clips):
            clips.append(image_clips[i])
        if i < len(video_clips):
            clips.append(video_clips[i])

    final_clip = concatenate_videoclips(clips, method="compose")

    # Load the overlay video
    overlay_video = VideoFileClip(overlay_video_path).resize(height=150).set_fps(frame_rate)

    # Apply any color corrections if needed
    overlay_video = overlay_video.fx(vfx.colorx, 1.0)  # Adjust colors if needed (1.0 means no change)

    # Create a circular mask using NumPy
    radius = overlay_video.h // 2
    circle_mask = np.zeros((overlay_video.h, overlay_video.w), dtype=np.uint8)
    y, x = np.ogrid[:overlay_video.h, :overlay_video.w]
    mask_center = (overlay_video.w // 2, overlay_video.h // 2)
    mask_area = (x - mask_center[0])*2 + (y - mask_center[1])*2 <= radius*2
    circle_mask[mask_area] = 255

    # Apply the circular mask to the overlay video
    # overlay_video = overlay_video.set_mask(ImageClip(circle_mask, ismask=True).set_duration(overlay_video.duration))

    # Loop the overlay video to match the length of the final clip
    overlay_video = overlay_video.loop(duration=final_clip.duration)

    # Position the overlay video at the bottom-right corner
    overlay_position = (common_resolution[0] - overlay_video.w - 10, common_resolution[1] - overlay_video.h - 10)  # 10px padding
    overlay_video = overlay_video.set_position(overlay_position)

    # Overlay the video on top of the slideshow
    final_composite = CompositeVideoClip([final_clip, overlay_video])

    # Load the audio file
    audio_clip = AudioFileClip(audio_path)
    final_composite = final_composite.set_audio(audio_clip)

    # Updated subtitle part
    subtitles = SubtitlesClip(
        srt_file_path, 
        lambda txt: TextClip(
            txt, 
            fontsize=10,        # Larger font size
            color='white',      # White text color
            bg_color='black',   # Black background color
            size=(final_composite.w*0.7, None)  # Width matching the video
        )
    )

    # Position subtitles at the bottom with some margin
    subtitles = subtitles.set_position(('center', 'bottom')).margin(bottom=20)

    # Add subtitles to the final video composite
    final_composite = CompositeVideoClip([final_composite, subtitles])

    try:
        final_composite.write_videofile(output_video_path, codec="libx264", audio_codec="aac")
    except Exception as e:
        print(f"Error creating slideshow video: {e}")
language_map = {
    'hindi': 'hi',
    'english': 'en',
    'bengali': 'bn',
    'telugu': 'te',
    'marathi': 'mr',
    'tamil': 'ta',
    'kannada': 'kn',
    'malayalam': 'ml',
    'gujarati': 'gu',
    'punjabi': 'pa',
    'urdu': 'ur'
}

# Initialize the Translator object
translator = Translator()

def translate_speech(text, target_language_code):
    try:
        # Translate the speech to the target language
        translation = translator.translate(text, dest=target_language_code)
        return translation.text
    except Exception as e:
        print(f"Error translating text to {target_language_code}: {e}")
        return text  # Return original text if translation fails



def generate_audio_from_text(text, output_audio_path, language_code='en'):
    try:
        # Generate the audio in the desired language
        tts = gTTS(text, lang=language_code)
        tts.save(output_audio_path)
        print(f"Audio saved to {output_audio_path} in {language_code}")
    except Exception as e:
        print(f"Error generating audio: {e}")

def speed_up_audio(input_audio_path, output_audio_path, background_music_path, speed=1.2, music_volume=-20):
    try:
        # Ensure the background music file exists
        if not os.path.exists(background_music_path):
            raise FileNotFoundError(f"Background music file not found: {background_music_path}")

        # Load the main audio and speed it up
        audio = AudioSegment.from_file(input_audio_path)
        sped_up_audio = audio.speedup(playback_speed=speed)

        # Load the background music
        background_music = AudioSegment.from_file(background_music_path)

        # Adjust the background music volume
        background_music = background_music - abs(music_volume)  # Lower the volume by music_volume dB

        # Loop the background music to match the length of the sped-up audio
        if len(background_music) < len(sped_up_audio):
            loop_count = len(sped_up_audio) // len(background_music) + 1
            background_music = background_music * loop_count
        background_music = background_music[:len(sped_up_audio)]

        # Mix the sped-up audio with the background music
        final_audio = sped_up_audio.overlay(background_music)

        # Export the final mixed audio
        final_audio.export(output_audio_path, format="mp3")
        print(f"Final audio with background music saved to {output_audio_path}")
    except Exception as e:
        print(f"Error generating final audio with background music: {e}")


if os.path.exists(audio_output_path):
    os.remove(audio_output_path)
    print(f"Deleted {audio_output_path}")