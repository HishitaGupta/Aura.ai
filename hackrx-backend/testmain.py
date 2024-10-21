import requests
import pdfplumber
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import os
import re
import google.generativeai as genai
from dotenv import load_dotenv
from transformers import pipeline
from moviepy.editor import VideoFileClip, ImageClip, concatenate_videoclips, CompositeVideoClip, TextClip, ColorClip, AudioFileClip
from moviepy.editor import *
from moviepy.video.fx.all import crop, loop
from moviepy.video.tools.drawing import circle
import json
import pysrt
import numpy as np
from gtts import gTTS
from googletrans import Translator
from pydub import AudioSegment
from moviepy.video.fx.all import colorx
from moviepy.video.tools.subtitles import SubtitlesClip
import base64
from docx import Document
from pptx import Presentation

# Set up folders
pictures_folder = "pictures"
videos_folder = "videos"
output_video_path = "final_slideshow.mp4"
audio_output_path = "final_audio.mp3"
audio_output_speedup_path = "final_audio_speeded_up.mp3"
srt_file_path = "subtitles.srt"

os.makedirs(pictures_folder, exist_ok=True)
os.makedirs(videos_folder, exist_ok=True)

# API URLs
img_url = 'https://pixabay.com/api/'
vid_url = 'https://pixabay.com/api/videos/'
STABILITY_API_URL = 'https://api.stability.ai/v2beta/stable-image/generate/core'
STABILITY_API_KEY = os.getenv("STABILITY_API_KEY")

# Load environment variables from .env file
load_dotenv()

API_KEY = os.getenv("IMG_API")
api_key = os.getenv("API_KEY")

# Configure API key for Google Gemini
genai.configure(api_key=api_key)
model = genai.GenerativeModel('gemini-1.5-flash')

# Set path for Tesseract OCR
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Define common resolution and frame rate
common_resolution = (1280, 720)
frame_rate = 24

def summarize_text(text):
    print("summarize_text TRIGGERED")
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    max_chunk = 1024
    text_chunks = [text[i:i + max_chunk] for i in range(0, len(text), max_chunk)]
    summary = ""
    for chunk in text_chunks:
        summary += summarizer(chunk, max_length=130, min_length=30, do_sample=False)[0]['summary_text'] + " "
    return summary.strip()

def extract_images_from_pdf(pdf_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    doc = fitz.open(pdf_path)
    image_list = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_list.append(image_bytes)
            image_path = os.path.join(output_folder, f"page_{page_num}_img_{img_index}.png")
            with open(image_path, "wb") as img_file:
                img_file.write(image_bytes)
    
    return image_list

def ocr_images_in_folder(folder_path):
    text = ""
    for filename in os.listdir(folder_path):
        if filename.endswith(".png"):
            image_path = os.path.join(folder_path, filename)
            try:
                image = Image.open(image_path)
                text += pytesseract.image_to_string(image) + "\n"
            except Exception as e:
                print(f"Error OCR processing {filename}: {e}")
    return text

def extract_text_from_pdf_images(pdf_path, output_folder):
    print("extract_text_from_pdf_images TRIGGERED")
    extract_images_from_pdf(pdf_path, output_folder)
    text = ocr_images_in_folder(output_folder)
    return text


def extract_text_from_txt(txt_path):
    with open(txt_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_pdf(pdf_path):
    print("extract_text_from_pdf TRIGGERED")
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text


def extract_text_from_file(file_path):
    if file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.txt'):
        return extract_text_from_txt(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file format")

def extract_chunks_from_srt(srt_file_path):
    subs = pysrt.open(srt_file_path)
    chunks = [sub.text for sub in subs]
    return chunks



def generate_keywords_from_summary(summary, srt_file_path):
    speeches = ""
    keywords = []

    # Generate speech from summary
    try:
        inp = model.generate_content(f"Write a minimum exact 250 and maximum exact 300 words more accurate summary based on the previous summary: {summary}.It should be plain text with no bullet points, in paragraphs, stick to the content, no '/n', (no CIN, UIN and GST details), and no bold stuff.")
        speeches = inp.text
        with open("speech_output.txt", 'w', encoding='utf-8') as output_file:
            output_file.write(speeches)
        print(f"Speech generated successfully: {output_txt_path}")
    except Exception as e:
        print(f"Error generating speech: {e}")

    # Extract chunks from SRT and generate keywords
    try:
        chunks = extract_chunks_from_srt(srt_file_path)

        for chunk in chunks:
            try:
                # Generate keywords for each chunk
                res = model.generate_content(f"Extract a list of key terms from the following chunk: {chunk}. The list should include important nouns (like proper nouns or objects) and action verbs (like payable, approval) that are important for visual representation and audio synchronization. For each extracted word, provide a two-word relatable visual metaphor suitable for depiction in images or videos. Exclude words that can't fetch relatable imagery.Example Sentence: 'Premium payable on renewal and on subsequent continuation of Bajaj Allianz are subject to change with prior approval from IRDA'.Expected Output:  Premium – Gold Coins,Payment – Credit Card,Renewal – Rotating Clock,Bajaj Allianz – Company Logo,Change – Circular Arrows,Approval – Green Checkmark,IRDA – Corporate Building.")
                keywords += re.findall(r'\*\*(.*?)\*\*', res.text)
            except Exception as e:
                print(f"Error generating keywords for chunk: {chunk}, Error: {e}")

        # Deduplicate keywords
        keywords = list(set(keywords))

        # Print generated keywords
        print(f"Generated Keywords: {keywords}")
    except Exception as e:
        print(f"Error extracting keywords from SRT: {e}")

    return {
        "speech": speeches,
        "keywords": keywords,
    }




# Subtitle part

def generate_subtitles_from_speech(speech_text, audio_duration, output_srt_path, chunk_duration=5):
    # Split speech into chunks
    words = speech_text.split()
    chunk_size = int(audio_duration // chunk_duration)
    chunks = [words[i:i + chunk_size] for i in range(0, len(words), chunk_size)]
    
    # Calculate total number of chunks
    total_chunks = len(chunks)
    
    # Adjust chunk_duration if it exceeds the audio duration
    if total_chunks > 1 and (audio_duration / total_chunks) < chunk_duration:
        chunk_duration = audio_duration / total_chunks
    
    # Create subtitles
    subs = pysrt.SubRipFile()
    for i, chunk in enumerate(chunks):
        start_time = i * (audio_duration / total_chunks)
        end_time = (i + 1) * (audio_duration / total_chunks)
        
        # Ensure end time does not exceed audio duration
        if end_time > audio_duration:
            end_time = audio_duration
        
        start_time = pysrt.SubRipTime(seconds=start_time)
        end_time = pysrt.SubRipTime(seconds=end_time)

        subtitle = pysrt.SubRipItem(index=i+1, start=start_time, end=end_time, text=' '.join(chunk))
        subs.append(subtitle)

    # Save the subtitles to a file
    subs.save(output_srt_path, encoding='utf-8')
    print(f"Subtitles saved to {output_srt_path}")

def generate_prompts_from_json(json_file_path, output_json_path):
    # Step 1: Read the JSON file
    try:
        with open(json_file_path, 'r', encoding='utf-8') as json_file:
            json_data = json.load(json_file)
    except Exception as e:
        print(f"Error reading JSON file: {e}")
        return []
    
    prompts = []
    
    # Step 2: Iterate over the JSON data
    for entry in json_data:
        # Extract the sentence (value) from the dictionary
        try:
            for idx, sentence in entry.items():
                print("Sentence: ", sentence)
                # Step 3: Generate a prompt for each sentence
                try:
                    inp = model.generate_content(
                        f"""Extract a list of key terms from the following text: {sentence} in the order they appear. The list should include important nouns (like proper nouns or objects) and action verbs (like payable, approval) that are important for visual representation and audio synchronization. For each extracted word, provide a two-word relatable visual metaphor suitable for depiction in images or videos. Exclude words that can't fetch relatable imagery. Example Sentence: 'Premium payable on renewal and on subsequent continuation of Bajaj Allianz are subject to change with prior approval from IRDA.Expected Output:Premium – Gold Coins,Payment – Credit Card,Renewal – Rotating Clock,Bajaj Allianz – Company Logo,Change – Circular Arrows,Approval – Green Checkmark,IRDA – Corporate Building.
                        make sure there's no whitespace in visual metaphors, if there's whitespace then consider the word after whitespace
                        return json data in list format of keyword and visual metaphor pair.
                        ```json
                        [
                            {{"keyword": "Premium", "visual_metaphor": "Coins"}},
                            {{"keyword": "Payment", "visual_metaphor": "Card"}},
                            {{"keyword": "Renewal", "visual_metaphor": "Clock"}},
                            {{"keyword": "Bajaj Allianz", "visual_metaphor": "Logo"}},
                            {{"keyword": "Change", "visual_metaphor": "Arrows"}},
                            {{"keyword": "Approval", "visual_metaphor": "Checkmark"}},
                            {{"keyword": "IRDA", "visual_metaphor": "Building"}}
                        ]
                        ```
                        """
                    )
                    
                    # Step 4: Remove the '```json' and '```' from the response to isolate the JSON content
                    json_extracted = inp.text.replace('```json', '').replace('```', '').strip()
                    
                    # Step 5: Convert the cleaned string into a valid JSON object
                    try:
                        json_obj = json.loads(json_extracted)
                        prompts.append(json_obj)
                    except json.JSONDecodeError as decode_error:
                        print(f"Error decoding JSON for sentence: {sentence}, Error: {decode_error}")
                except Exception as gen_error:
                    print(f"Error generating prompt for sentence: {sentence}, Error: {gen_error}")
        except Exception as e:
            print(f"Error processing JSON entry: {entry}, Error: {e}")
    
    # Step 6: Write the resulting prompts to the output JSON file
    try:
        with open(output_json_path, 'w', encoding='utf-8') as output_file:
            json.dump(prompts, output_file, ensure_ascii=False, indent=4)
        print(f"Prompts successfully saved to {output_json_path}")
    except Exception as file_error:
        print(f"Error writing to JSON file: {file_error}")
    
    return prompts

generate_prompts_from_json(r"C:\Users\Happy yadav\Desktop\Technology\hack\hackrx-backend\speech_output.json", "final-prompt.json")

# def create_audio_and_get_word_durations(json_file_path, output_folder):
#     # Step 1: Create the output folder if it doesn't exist
#     os.makedirs(output_folder, exist_ok=True)

#     # Step 2: Read the JSON file
#     try:
#         with open(json_file_path, 'r', encoding='utf-8') as json_file:
#             json_data = json.load(json_file)  # json_data is a list of lists
#     except Exception as e:
#         print(f"Error reading JSON file: {e}")
#         return

#     # Step 3: Iterate through each list in the JSON data
#     for entry_list in json_data:  # entry_list is a list of dictionaries
#         for entry in entry_list:  # entry is a dictionary
#             sentence = entry.get("keyword", "")  # Extract the sentence
#             print(f"Generating audio for sentence: {sentence}")
#             try:
#                 # Generate audio using gTTS
#                 tts = gTTS(text=sentence, lang='en')
#                 audio_file_path = os.path.join(output_folder, f"audio_{sentence.replace(' ', '_')}.mp3")
                
#                 # Save the audio file
#                 tts.save(audio_file_path)
#                 print(f"Saved audio file: {audio_file_path}")

#                 # Step 4: Analyze the audio to get word durations
#                 word_durations = get_word_durations(sentence, audio_file_path)
#                 print(f"Word durations for sentence '{sentence}': {word_durations}")

#             except Exception as audio_error:
#                 print(f"Error generating audio for sentence: {sentence}, Error: {audio_error}")

# # Call the function
# create_audio_and_get_word_durations("final-prompt.json", "audio_files")


# def create_single_audio_from_json(json_file_path, output_file_path):
#     # Step 1: Create the output folder if it doesn't exist
#     output_folder = os.path.dirname(output_file_path)
#     os.makedirs(output_folder, exist_ok=True)

#     # Step 2: Read the JSON file
#     try:
#         with open(json_file_path, 'r', encoding='utf-8') as json_file:
#             json_data = json.load(json_file)  # json_data is a list of dictionaries
#     except Exception as e:
#         print(f"Error reading JSON file: {e}")
#         return

#     # Initialize an empty AudioSegment to hold the final combined audio
#     combined_audio = AudioSegment.silent(duration=0)  # Start with silence

#     # Step 3: Iterate through each dictionary in the JSON data
#     for entry in json_data:  # entry is a dictionary, e.g., {"0": "A Non-Medical Expenses Cover..."}
#         # Extract the sentence from the dictionary (there's only one key-value pair per dictionary)
#         sentence = next(iter(entry.values()))  # Get the sentence value
#         print(f"Generating audio for sentence: {sentence}")
#         try:
#             # Generate audio using gTTS
#             tts = gTTS(text=sentence, lang='en')
#             temp_audio_path = os.path.join(output_folder, f"temp_audio.mp3")

#             # Save the audio file
#             tts.save(temp_audio_path)
#             print(f"Saved audio file: {temp_audio_path}")

#             # Step 4: Load the audio file into an AudioSegment
#             audio_segment = AudioSegment.from_mp3(temp_audio_path)

#             # Combine the current audio segment with the overall combined audio
#             combined_audio += audio_segment + AudioSegment.silent(duration=500)  # Add 500ms silence between sentences

#         except Exception as audio_error:
#             print(f"Error generating audio for sentence: {sentence}, Error: {audio_error}")

#     # Step 5: Save the combined audio to a final file
#     combined_audio.export(output_file_path, format='mp3')
#     print(f"Combined audio file saved as: {output_file_path}")

# # Call the function to create a single audio file from the JSON
# create_single_audio_from_json("speech_output.json", "audio_files/combined_audio.mp3")

def create_single_and_individual_audios_from_json(json_file_path, output_folder):
    # Step 1: Create the output folder if it doesn't exist
    os.makedirs(output_folder, exist_ok=True)

    # Step 2: Read the JSON file
    try:
        with open(json_file_path, 'r', encoding='utf-8') as json_file:
            json_data = json.load(json_file)  # json_data is a list of dictionaries
    except Exception as e:
        print(f"Error reading JSON file: {e}")
        return

    # Initialize an empty AudioSegment to hold the final combined audio
    combined_audio = AudioSegment.silent(duration=0)  # Start with silence
    audio_durations = []  # To store the audio file name and duration

    # Step 3: Iterate through each dictionary in the JSON data
    for index, entry in enumerate(json_data):  # entry is a dictionary, e.g., {"0": "A Non-Medical Expenses Cover..."}
        # Extract the sentence from the dictionary (there's only one key-value pair per dictionary)
        sentence = next(iter(entry.values()))  # Get the sentence value
        print(f"Generating audio for sentence: {sentence}")
        
        try:
            # Generate audio using gTTS
            tts = gTTS(text=sentence, lang='en')
            individual_audio_filename = f"audio_sentence_{index + 1}.mp3"
            individual_audio_path = os.path.join(output_folder, individual_audio_filename)

            # Save the individual audio file for the current sentence
            tts.save(individual_audio_path)
            print(f"Saved individual audio file: {individual_audio_path}")

            # Step 4: Load the individual audio file into an AudioSegment
            audio_segment = AudioSegment.from_mp3(individual_audio_path)

            # Get the duration of the audio segment in seconds
            audio_duration = len(audio_segment) // 1000  # Convert from milliseconds to seconds

            # Add the audio file name and duration to the list
            audio_durations.append({individual_audio_filename: str(audio_duration)})

            # Combine the current audio segment with the overall combined audio
            combined_audio += audio_segment + AudioSegment.silent(duration=500)  # Add 500ms silence between sentences

        except Exception as audio_error:
            print(f"Error generating audio for sentence: {sentence}, Error: {audio_error}")

    # Step 5: Save the combined audio to a final file
    combined_audio_path = os.path.join(output_folder, "combined_audio.mp3")
    combined_audio.export(combined_audio_path, format='mp3')
    print(f"Combined audio file saved as: {combined_audio_path}")

    # Step 6: Write the audio durations to a JSON file
    durations_json_path = os.path.join(output_folder, "audio_durations.json")
    try:
        with open(durations_json_path, 'w', encoding='utf-8') as json_file:
            json.dump(audio_durations, json_file, ensure_ascii=False, indent=4)
        print(f"Audio durations saved as: {durations_json_path}")
    except Exception as json_error:
        print(f"Error writing audio durations to JSON file: {json_error}")

# Call the function to create both individual and combined audio files
create_single_and_individual_audios_from_json("speech_output.json", "audio_files")

def get_word_durations(sentence, audio_file_path):
    # Step 1: Load the audio file
    audio = AudioSegment.from_mp3(audio_file_path)

    # Step 2: Split the sentence into words
    words = sentence.split()
    durations = {}
    
    # Step 3: Calculate the duration for each word
    total_length = len(audio)
    word_length = len(words)
    average_duration_per_word = total_length / word_length
    
    # Create a list of words with their durations
    for word in words:
        durations[word] = average_duration_per_word

    return durations

def calculate_keyword_timestamps(speech_output_path, audio_durations_path, combined_audio_path, final_prompts_path):
    # Load audio duration of the combined audio file
    try:
        combined_audio = AudioSegment.from_mp3(combined_audio_path)
        total_duration = len(combined_audio) // 1000  # duration in seconds
    except Exception as e:
        print(f"Error reading combined audio file: {e}")
        return
    
    # Load speech output from text file
    try:
        with open(speech_output_path, 'r', encoding='utf-8') as text_file:
            speech_data = text_file.readlines()  # Read each line as a sentence
    except Exception as e:
        print(f"Error reading speech output text file: {e}")
        return

    # Calculate total words
    total_words = sum(len(sentence.split()) for sentence in speech_data)

    # Calculate average duration per word
    avg_duration_per_word = total_duration / total_words if total_words > 0 else 0

    # Load final prompts (keywords)
    try:
        with open(final_prompts_path, 'r', encoding='utf-8') as json_file:
            final_prompts = json.load(json_file)
    except Exception as e:
        print(f"Error reading final prompts JSON: {e}")
        return

    keyword_timestamps = {}

    # Iterate through each sentence to find keywords
    current_time = 0
    for sentence in speech_data:
        words = sentence.split()
        
        for keyword_group in final_prompts:
            for keyword_data in keyword_group:
                keyword = keyword_data["keyword"]
                # Check if the keyword is in the current sentence
                if keyword in sentence:
                    # Find the position of the keyword
                    word_index = words.index(keyword) if keyword in words else -1
                    if word_index != -1:
                        # Calculate the timestamp for this keyword
                        keyword_timestamp = current_time + (word_index * avg_duration_per_word)
                        
                        # Store the timestamp for each unique keyword
                        if keyword not in keyword_timestamps:
                            keyword_timestamps[keyword] = round(keyword_timestamp, 2)  # Round to 2 decimal places

        # Update current time for the next sentence
        current_time += len(words) * avg_duration_per_word  # Update current time based on word count in the sentence

    # Output keyword timestamps
    output_path = os.path.join("audio_files", "keyword_timestamps.json")
    try:
        with open(output_path, 'w', encoding='utf-8') as json_file:
            json.dump(keyword_timestamps, json_file, ensure_ascii=False, indent=4)
        print(f"Keyword timestamps saved to: {output_path}")
    except Exception as e:
        print(f"Error saving keyword timestamps: {e}")

calculate_keyword_timestamps("speech_output.json", "audio_files/audio_durations.json", "audio_files/combined_audio.mp3","final-prompt.json")



# get_word_durations()

# Quiz part

def generate_quiz(text):
    quiz_string = ""
    try:
        inp = model.generate_content(f"Generate quiz which contains 3 questions with unique answers in MCQ format containing 'question', 'options', 'answer' in JSON list format based on the text: {text}")
        print(inp.text)
        quiz_string = inp.text
    except Exception as e:
        print(f"Error generating quiz: {e}")
    return quiz_string

def save_quiz_to_json(quiz_string, output_file):
    cleaned_string = quiz_string.replace('```json', '').replace('```', '').strip()
    
    try:
        quiz_json = json.loads(cleaned_string)
        with open(output_file, 'w') as json_file:
            json.dump(quiz_json, json_file, indent=4)
        print(f"Quiz data saved successfully to {output_file}")
    except json.JSONDecodeError as e:
        print(f"Failed to decode JSON: {e}")

def clean_text(text):
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'[^a-zA-Z0-9\s.,]', '', text)
    return text.strip()

def get_text_from_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as text_file:
            text = text_file.read()
        return text
    except FileNotFoundError:
        print(f"The file {file_path} does not exist.")
        return None
    except Exception as e:
        print(f"An error occurred while reading the file: {e}")
        return None

def ask_aura_question(user_question, txt_file_path):
    answer = ""
    try:
        # Get the extracted text from the txt file
        extracted_text = get_text_from_txt(txt_file_path)
        
        if not extracted_text:
            return "Error: Could not retrieve text from file."

        # Use the extracted text and the user's question to prompt Gemini AI
        prompt = f"Answer the following question based on this text: '{extracted_text}'.\nQuestion: {user_question}"
        response = model.generate_content(prompt)
        print("Answer response: ", response.text)  # Debug print for checking the response
        answer = response.text
    except Exception as e:
        print(f"Error while asking Gemini AI: {e}")
        return f"Error: {e}"
    
    return answer

def save_video_from_url(video_url, save_directory, keyword):
    """Download and save video from the specified URL with a unique index."""
    os.makedirs(save_directory, exist_ok=True)
    
    try:
        video_response = requests.get(video_url)
        if video_response.status_code == 200:
            # Use the keyword name in the file name to avoid duplication
            video_path = os.path.join(save_directory, f'{keyword}.mp4')
            with open(video_path, 'wb') as file:
                file.write(video_response.content)
            print(f"Video saved as {video_path}")
            return True  # Successfully saved one video
        else:
            print(f"Failed to download video. Status code: {video_response.status_code}")
            return False
    except Exception as e:
        print(f"Error saving video: {e}")
        return False

def generate_and_save_videos_for_visual_metaphors(final_prompts_path):
    """Fetch and save one video per visual metaphor from the final prompts JSON."""
    
    # Load final prompts (visual metaphors)
    try:
        with open(final_prompts_path, 'r', encoding='utf-8') as json_file:
            final_prompts = json.load(json_file)
    except Exception as e:
        print(f"Error reading final prompts JSON: {e}")
        return

    # Ensure the videos folder exists
    os.makedirs(videos_folder, exist_ok=True)

    # Track saved visual metaphors to avoid duplicates
    saved_metaphors = set()

    # Loop through visual metaphors to fetch only one video per metaphor
    for group_index, keyword_group in enumerate(final_prompts):
        for keyword_data in keyword_group:
            visual_metaphor = keyword_data["visual_metaphor"].replace(" ", "_")  # Replace spaces for file naming
            keyword = keyword_data["keyword"]
            # Skip if this visual metaphor has already been saved
            if visual_metaphor in saved_metaphors:
                print(f"Skipping already saved visual metaphor: {visual_metaphor}")
                continue
            
            print(f"Processing visual metaphor {group_index + 1}: {visual_metaphor}")

            params = {
                'key': API_KEY,
                'q': visual_metaphor,
                'video_type': 'all',  # For videos
                'per_page': 3         # Request up to 3 results (in case some are missing)
            }

            try:
                # Fetch videos from Pixabay
                response = requests.get(vid_url, params=params)
                if response.status_code == 200:
                    data = response.json()
                    videos = data.get('hits', [])
                    if videos:
                        for idx, video in enumerate(videos):
                            video_url = video.get('videos', {}).get('medium', {}).get('url', '')
                            if video_url:
                                # Save the video and add the metaphor to the saved set
                                if save_video_from_url(video_url, videos_folder, keyword):
                                    saved_metaphors.add(visual_metaphor)  # Mark as saved
                                    break  # Exit loop after saving one video
                    else:
                        print(f"No videos found for '{visual_metaphor}'")
                else:
                    print(f"Failed to fetch videos for '{visual_metaphor}'. Status code: {response.status_code}")

            except Exception as e:
                print(f"Error processing visual metaphor '{visual_metaphor}': {e}")

# Example usage
generate_and_save_videos_for_visual_metaphors("final-prompt.json")

def trim_video(video_path, duration=7):
    try:
        with VideoFileClip(video_path) as video:
            print(f"Original Duration: {video.duration}")
            if video.duration > duration:
                trimmed_video = video.subclip(0, duration)
                trimmed_video_path = video_path.replace(".mp4", "_trimmed.mp4")
                trimmed_video.write_videofile(trimmed_video_path, codec='libx264', audio_codec='aac')
                print(f"Trimmed Duration: {trimmed_video.duration}")
                os.replace(trimmed_video_path, video_path)
    except Exception as e:
        print(f"Error trimming video {video_path}: {e}")
        if os.path.exists(video_path):
            os.remove(video_path)

def generate_image_from_prompt(prompt, images_folder, image_index):
    headers = {
        'Authorization': f'Bearer sk-Xips2HbQGXa2GTbRE3Y08s933bqp7eNhSUGw4eyFW1Vu4Cep',
        
        'accept': 'image/*'
    }

    data = {
        "prompt": prompt,
        # "aspect_ratio": "16:9",
        "output_format": "png",
        "aspect_ratio": "16:9",
    }
    # print("data: ", data['prompt'])

    files = {
        'prompt': (None, prompt)  # Use None to indicate that it's not a file
    }

    print("prompt: ", files['prompt'])
    try:
        print("Sending prompt:", prompt)
        response = requests.post(
            "https://api.stability.ai/v2beta/stable-image/generate/core",
            headers=headers,
            files=files,  # Sending as multipart/form-data
            data=data
        )

        # Check the response status
        print("Response Status Code:", response.status_code)
        # print("Response Content:", response.text)

        response.raise_for_status()  # Raise an error for bad responses

        # If receiving the image directly
        if headers['accept'] == 'image/*':
            image_data = response.content
            image_path = os.path.join(images_folder, f"generated_image_{image_index}.png")
            with open(image_path, "wb") as img_file:
                img_file.write(image_data)
            print(f"Image saved: {image_path}")
        else:
            # If receiving base64 JSON response
            json_response = response.json()
            base64_image = json_response.get('image', None)
            if base64_image:
                image_data = base64.b64decode(base64_image)
                image_path = os.path.join(images_folder, f"generated_image_{image_index}.png")
                with open(image_path, "wb") as img_file:
                    img_file.write(image_data)
                print(f"Image saved: {image_path}")
            else:
                print("No image data received in JSON response.")

    except requests.exceptions.HTTPError as http_err:
        print(f"HTTP error occurred: {http_err}")
    except Exception as e:
        print(f"Error generating image from prompt: {e}")


#pixaby code 

def generate_and_save_images_and_videos_for_keywords(keywords):
    # Ensure the headers are set up correctly
    headers = {'Authorization': f'Bearer {API_KEY}'}
    
    # Loop through keywords to fetch images and videos
    for i, keyword in enumerate(keywords):
        print(f"Processing keyword {i + 1}: {keyword}")

        params = {
            'key': API_KEY,
            'q': keyword,
            'video_type': 'all',    # For videos
            'per_page': 3            # Number of results per page
        }

        try:
            if i < 10:  # Fetch videos
                response = requests.get(vid_url, params=params)
                if response.status_code == 200:
                    data = response.json()
                    videos = data.get('hits', [])
                    for idx, video in enumerate(videos):
                        video_url = video.get('videos', {}).get('medium', {}).get('url', '')
                        if video_url:
                            save_video_from_url(video_url, videos_folder, i + 1 + idx)
                            video_path = os.path.join(videos_folder, f'video_{i + 1 + idx}.mp4')
                            trim_video(video_path)
                else:
                    print(f"Failed to fetch videos for keyword {i + 1}. Status code: {response.status_code}")

            else:
                print("Nothing is fetching...")
        
        except Exception as e:
            print(f"Error processing keyword {keyword}: {e}")

def clean_up_videos(folder_path, max_duration=7):
    for filename in os.listdir(folder_path):
        if filename.endswith(".mp4"):
            video_path = os.path.join(folder_path, filename)
            try:
                with VideoFileClip(video_path) as video:
                    if video.duration > max_duration:
                        os.remove(video_path)
                        print(f"Deleted {video_path} because its duration was greater than {max_duration} seconds.")
            except Exception as e:
                print(f"Error checking video duration for {video_path}: {e}")

def create_speech_file(output_txt_path, output_json_path):
    try:
        # Step 1: Read from the speech_output.txt file
        with open(output_txt_path, 'r', encoding='utf-8') as txt_file:
            speech_text = txt_file.read()
        
        # Step 2: Split the text into sentences based on periods (dots)
        sentences = [sentence.strip() for sentence in speech_text.split('.') if sentence.strip()]
        
        # Step 3: Create the JSON data by indexing the sentences
        json_output = [{i: sentence} for i, sentence in enumerate(sentences)]
        
        # Step 4: Save the JSON data into the specified output JSON file
        with open(output_json_path, 'w', encoding='utf-8') as json_file:
            json.dump(json_output, json_file, ensure_ascii=False, indent=4)
        
        print(f"JSON file created successfully: {output_json_path}")
    
    except Exception as e:
        print(f"An error occurred: {e}")

output_txt_path = 'speech_output.txt'
output_json_path = 'speech_output.json'
create_speech_file(output_txt_path, output_json_path)


def create_slideshow_with_audio(images_folder, videos_folder, output_video_path, audio_path, overlay_video_path, srt_file_path, image_duration=2, fade_duration=1):
    image_clips=[]
    video_clips = []

    # Calculate the total duration of the audio
    audio_clip = AudioFileClip(audio_path)
    audio_duration = audio_clip.duration
    audio_clip.close()

    # Calculate the number of images and videos
    num_images = len([f for f in os.listdir(images_folder) if f.endswith(('.jpg', '.png'))])
    num_videos = len([f for f in os.listdir(videos_folder) if f.endswith('.mp4')])

    # Calculate the duration per image and video to match the audio length
    if  num_images + num_videos > 0:
        duration_per_clip = audio_duration / (num_images+num_videos)

    # Load and process image clips
    for filename in sorted(os.listdir(images_folder)):
        if filename.endswith(('.jpg', '.png')):
            image_path = os.path.join(images_folder, filename)
            image_clip = ImageClip(image_path, duration=duration_per_clip).set_fps(frame_rate)
            
            # Resize while maintaining aspect ratio
            image_clip = image_clip.resize(height=common_resolution[1])
            image_clip = image_clip.set_duration(duration_per_clip).fadein(fade_duration).fadeout(fade_duration)

            background = ColorClip(size=common_resolution, color=(0, 0, 0), duration=duration_per_clip)
            image_clip = CompositeVideoClip([background, image_clip.set_position("center")])
            image_clips.append(image_clip)

    # Load and process video clips
    for filename in sorted(os.listdir(videos_folder)):
        if filename.endswith('.mp4'):
            video_path = os.path.join(videos_folder, filename)
            video_clip = VideoFileClip(video_path).set_fps(frame_rate)

            video_clip = video_clip.resize(height=common_resolution[1])
            video_clip = video_clip.set_duration(duration_per_clip).fadein(fade_duration).fadeout(fade_duration)

            background = ColorClip(size=common_resolution, color=(0, 0, 0), duration=duration_per_clip)
            video_clip = CompositeVideoClip([background, video_clip.set_position("center")])
            video_clips.append(video_clip)

    # Create the final list of clips by alternating between images and videos
    clips = []
    max_len = max(len(image_clips), len(video_clips))
    # max_len = len(video_clips)

    for i in range(max_len):
        if i < len(image_clips):
            clips.append(image_clips[i])
        if i < len(video_clips):
            clips.append(video_clips[i])

    final_clip = concatenate_videoclips(clips, method="compose")

    # Load the overlay video
    overlay_video = VideoFileClip(overlay_video_path).resize(height=150).set_fps(frame_rate)

    # Apply any color corrections if needed
    overlay_video = overlay_video.fx(vfx.colorx, 1.0)  # Adjust colors if needed (1.0 means no change)

    # Create a circular mask using NumPy
    radius = overlay_video.h // 2
    circle_mask = np.zeros((overlay_video.h, overlay_video.w), dtype=np.uint8)
    y, x = np.ogrid[:overlay_video.h, :overlay_video.w]
    mask_center = (overlay_video.w // 2, overlay_video.h // 2)
    mask_area = (x - mask_center[0])**2 + (y - mask_center[1])**2 <= radius**2
    circle_mask[mask_area] = 255

    # Apply the circular mask to the overlay video
    # overlay_video = overlay_video.set_mask(ImageClip(circle_mask, ismask=True).set_duration(overlay_video.duration))

    # Loop the overlay video to match the length of the final clip
    overlay_video = overlay_video.loop(duration=final_clip.duration)

    # Position the overlay video at the bottom-right corner
    overlay_position = (common_resolution[0] - overlay_video.w - 10, common_resolution[1] - overlay_video.h - 10)  # 10px padding
    overlay_video = overlay_video.set_position(overlay_position)

    # Overlay the video on top of the slideshow
    final_composite = CompositeVideoClip([final_clip, overlay_video])

    # Load the audio file
    audio_clip = AudioFileClip(audio_path)
    final_composite = final_composite.set_audio(audio_clip)

    # Updated subtitle part
    subtitles = SubtitlesClip(
        srt_file_path, 
        lambda txt: TextClip(
            txt, 
            fontsize=10,        # Larger font size
            color='white',      # White text color
            bg_color='black',   # Black background color
            size=(final_composite.w*0.7, None)  # Width matching the video
        )
    )

    # Position subtitles at the bottom with some margin
    subtitles = subtitles.set_position(('center', 'bottom')).margin(bottom=20)

    # Add subtitles to the final video composite
    final_composite = CompositeVideoClip([final_composite, subtitles])

    try:
        final_composite.write_videofile(output_video_path, codec="libx264", audio_codec="aac")
    except Exception as e:
        print(f"Error creating slideshow video: {e}")
language_map = {
    'hindi': 'hi',
    'english': 'en',
    'bengali': 'bn',
    'telugu': 'te',
    'marathi': 'mr',
    'tamil': 'ta',
    'kannada': 'kn',
    'malayalam': 'ml',
    'gujarati': 'gu',
    'punjabi': 'pa',
    'urdu': 'ur'
}

# Initialize the Translator object
translator = Translator()

def translate_speech(text, target_language_code):
    try:
        # Translate the speech to the target language
        translation = translator.translate(text, dest=target_language_code)
        return translation.text
    except Exception as e:
        print(f"Error translating text to {target_language_code}: {e}")
        return text  # Return original text if translation fails



def generate_audio_from_text(text, output_audio_path, language_code='en'):
    try:
        # Generate the audio in the desired language
        tts = gTTS(text, lang=language_code)
        tts.save(output_audio_path)
        print(f"Audio saved to {output_audio_path} in {language_code}")
    except Exception as e:
        print(f"Error generating audio: {e}")

def speed_up_audio(input_audio_path, output_audio_path, background_music_path, speed=1.2, music_volume=-20):
    try:
        # Ensure the background music file exists
        if not os.path.exists(background_music_path):
            raise FileNotFoundError(f"Background music file not found: {background_music_path}")

        # Load the main audio and speed it up
        audio = AudioSegment.from_file(input_audio_path)
        sped_up_audio = audio.speedup(playback_speed=speed)

        # Load the background music
        background_music = AudioSegment.from_file(background_music_path)

        # Adjust the background music volume
        background_music = background_music - abs(music_volume)  # Lower the volume by `music_volume` dB

        # Loop the background music to match the length of the sped-up audio
        if len(background_music) < len(sped_up_audio):
            loop_count = len(sped_up_audio) // len(background_music) + 1
            background_music = background_music * loop_count
        background_music = background_music[:len(sped_up_audio)]

        # Mix the sped-up audio with the background music
        final_audio = sped_up_audio.overlay(background_music)

        # Export the final mixed audio
        final_audio.export(output_audio_path, format="mp3")
        print(f"Final audio with background music saved to {output_audio_path}")
    except Exception as e:
        print(f"Error generating final audio with background music: {e}")


if os.path.exists(audio_output_path):
    os.remove(audio_output_path)
    print(f"Deleted {audio_output_path}")